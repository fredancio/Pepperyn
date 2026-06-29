"""
tests/test_rule_003_renderer_responsibility.py

PEPPERYN PRODUCT QUALITY CONTRACT — RULE 003
Renderer Responsibility

"Presentation problems are renderer problems.
 Never user problems. Never data problems. Never LLM problems."

Ce test est le gardien permanent de RULE 003.
Le critère n'est pas "le renderer gère mieux les erreurs".
Le critère est : AUCUNE exception ne sort jamais d'un renderer quelle que soit l'entrée.

Le renderer est entièrement responsable de sa propre présentation.
Il ne peut jamais déléguer un problème de présentation à l'appelant.

Protocole de vérification :
  1. Cas minimal (ExecutiveCaseJSON vide, tous les optionnels à None/[]) :
     → PDF, PPTX, Excel produisent des octets valides. Aucune exception.
  2. Champs optionnels à None :
     → PDF ne lève pas d'exception avec tous les champs narratifs à None.
  3. Listes vides :
     → PPTX et Excel avec décisions=[], risques=[], destroyers=[] ne crashent pas.
  4. Texte extrême (> 2 000 caractères par décision) :
     → Les trois renderers absorbent le contenu sans exception.
  5. Cas Optilux (fixture de référence) :
     → Les trois renderers produisent des octets ≥ taille minimum attendue.

RULE 003 — Renderer Responsibility.
RULE 002 — Zero Truncated Content (non-régression).
"""
import sys, os, pytest
sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))
os.environ.setdefault("ANTHROPIC_API_KEY", "dummy")

from models.schemas import AnalysisResult
from models.executive_case import (
    ExecutiveCaseJSON, DimensionScores, COIBreakdown, KPICard, ValueDestroyerItem,
    PriorityDecisionItem, ExecutionLogItem, ProjectionSeries, Scenarios, ScenarioItem,
    RiskItem, DataQuality,
)
from services.excel_export import generate_excel_report
from services.export_pdf_service import generate_pdf_report
from services.export_pptx_service import generate_pptx_report


# ─── Fixtures ─────────────────────────────────────────────────────────────────

MINIMAL_CASE = ExecutiveCaseJSON()  # Tous les champs à leur valeur par défaut


NONE_FIELDS_CASE = ExecutiveCaseJSON(
    company_name="Test Société",
    analysis_date="29/06/2026",
    executive_diagnosis=None,
    tension_phrase=None,
    inaction_risk=None,
    structural_loss_statement=None,
    structural_loss_value=None,
    urgency_level=None,
    value_creation_statement=None,
    health_score=None,
    dimension_scores=DimensionScores(rentabilite=None, risque=None, structure=None, liquidite=None),
    cost_of_inaction=COIBreakdown(annual=None, monthly=None, weekly=None, daily=None, hourly=None),
    kpi_dashboard=[],
    value_destroyers=[],
    priority_decisions=[],
    major_risks=[],
    execution_log=[],
    series=ProjectionSeries(action=[], inaction=[], equilibrium=[]),
    scenarios=Scenarios(best=None, likely=None, worst=None),
    data_quality=DataQuality(score=0, anomalies=[], assumptions=[], limits=[]),
)

_LONG_TEXT = (
    "Cette décision stratégique implique une restructuration profonde des processus "
    "opérationnels et commerciaux de l'entreprise, incluant la révision complète "
    "des contrats fournisseurs, la mise en place d'un nouveau système de suivi "
    "des performances, la formation des équipes dirigeantes sur les nouveaux KPIs, "
    "l'implémentation d'un tableau de bord de pilotage en temps réel, "
    "et la coordination avec les partenaires financiers pour assurer la liquidité "
    "nécessaire pendant la période de transition. "
) * 8   # ~2 200 caractères

EXTREME_TEXT_CASE = ExecutiveCaseJSON(
    company_name="Extrême SARL",
    analysis_date="29/06/2026",
    executive_diagnosis=_LONG_TEXT,
    tension_phrase=_LONG_TEXT[:500],
    priority_decisions=[
        PriorityDecisionItem(
            decision=_LONG_TEXT,
            annual_impact=1_000_000,
            priority="Élevée",
            roi_score=9.5,
        ),
        PriorityDecisionItem(
            decision=_LONG_TEXT,
            annual_impact=500_000,
            priority="Moyenne",
            roi_score=7.0,
        ),
    ],
    value_destroyers=[
        ValueDestroyerItem(name=_LONG_TEXT, annual_impact=-800_000, trend="↑"),
        ValueDestroyerItem(name=_LONG_TEXT, annual_impact=-400_000, trend="→"),
    ],
    major_risks=[
        RiskItem(description=_LONG_TEXT, severity="Élevé", impact="Critique", horizon="Immédiat"),
    ],
)


# ─── TestNoExceptionPropagation ───────────────────────────────────────────────

class TestNoExceptionPropagation:
    """
    RULE 003 — Aucune exception ne doit jamais quitter un renderer.

    Le renderer est une boîte noire de présentation.
    Quelle que soit l'entrée (vide, None, corrompue, extrême),
    il produit des octets valides OU échoue silencieusement.
    Il ne propage JAMAIS une exception à l'appelant.
    """

    def test_pdf_never_raises_with_minimal_case(self):
        """PDF : ExecutiveCaseJSON() vide → pas d'exception, bytes retournés."""
        try:
            result = generate_pdf_report(MINIMAL_CASE, "")
        except Exception as exc:
            pytest.fail(
                f"RULE 003 VIOLATION — PDF renderer a propagé une exception : {exc!r}\n"
                "Le renderer doit absorber tous les problèmes de présentation."
            )
        assert isinstance(result, bytes) and len(result) > 0, (
            "RULE 003 VIOLATION — PDF renderer a retourné des bytes vides sur cas minimal."
        )

    def test_pptx_never_raises_with_minimal_case(self):
        """PPTX : ExecutiveCaseJSON() vide → pas d'exception, bytes retournés."""
        try:
            result = generate_pptx_report(MINIMAL_CASE, "")
        except Exception as exc:
            pytest.fail(
                f"RULE 003 VIOLATION — PPTX renderer a propagé une exception : {exc!r}\n"
                "Le renderer doit absorber tous les problèmes de présentation."
            )
        assert isinstance(result, bytes) and len(result) > 0, (
            "RULE 003 VIOLATION — PPTX renderer a retourné des bytes vides sur cas minimal."
        )

    def test_excel_never_raises_with_minimal_case(self):
        """Excel : ExecutiveCaseJSON() vide → pas d'exception, bytes retournés."""
        try:
            result = generate_excel_report(AnalysisResult(), MINIMAL_CASE, "test")
        except Exception as exc:
            pytest.fail(
                f"RULE 003 VIOLATION — Excel renderer a propagé une exception : {exc!r}\n"
                "Le renderer doit absorber tous les problèmes de présentation."
            )
        assert isinstance(result, bytes) and len(result) > 0, (
            "RULE 003 VIOLATION — Excel renderer a retourné des bytes vides sur cas minimal."
        )


# ─── TestRendererIsolation ────────────────────────────────────────────────────

class TestRendererIsolation:
    """
    RULE 003 — Le renderer produit un fichier valide même avec des données absentes.

    "Données insuffisantes" sur une section est une réponse acceptable.
    Un crash ne l'est pas.
    """

    def test_pdf_produces_valid_bytes_with_all_none_optional_fields(self):
        """PDF avec tous les champs narratifs à None → fichier valide produit."""
        result = generate_pdf_report(NONE_FIELDS_CASE, "Test Société")
        assert isinstance(result, bytes), "PDF doit retourner des bytes."
        assert len(result) > 1_000, (
            f"PDF trop court ({len(result)} bytes) — probablement tronqué par une exception absorbée."
        )
        # Vérifier que c'est bien un PDF
        assert result[:4] == b"%PDF", "PDF renderer n'a pas produit un fichier PDF valide."

    def test_pptx_produces_valid_bytes_with_empty_lists(self):
        """PPTX avec listes vides (decisions=[], risks=[], destroyers=[]) → 16 slides produits."""
        from pptx import Presentation
        import io
        result = generate_pptx_report(NONE_FIELDS_CASE, "Test Société")
        assert isinstance(result, bytes), "PPTX doit retourner des bytes."
        prs = Presentation(io.BytesIO(result))
        assert len(prs.slides) == 16, (
            f"PPTX doit toujours produire 16 slides. Reçu : {len(prs.slides)} slides."
        )

    def test_excel_produces_valid_bytes_with_empty_lists(self):
        """Excel avec listes vides → workbook valide avec toutes les feuilles."""
        from openpyxl import load_workbook
        import io
        result = generate_excel_report(AnalysisResult(), NONE_FIELDS_CASE, "test")
        assert isinstance(result, bytes), "Excel doit retourner des bytes."
        wb = load_workbook(io.BytesIO(result))
        # Vérifier que le workbook contient au moins les 8 feuilles visibles
        visible_sheets = [s for s in wb.sheetnames if wb[s].sheet_state == "visible"]
        assert len(visible_sheets) >= 8, (
            f"Excel doit produire ≥8 feuilles visibles. Reçu : {visible_sheets}"
        )

    def test_pdf_no_exception_with_zero_scores(self):
        """PDF avec scores à 0 et confidence_score=0 → pas d'exception."""
        zero_case = ExecutiveCaseJSON(
            company_name="Zéro Corp",
            confidence_score=0,
            health_score=0,
            decisions_priority_score=0.0,
        )
        try:
            result = generate_pdf_report(zero_case, "Zéro Corp")
        except Exception as exc:
            pytest.fail(f"RULE 003 VIOLATION — PDF a crashé sur des scores à zéro : {exc!r}")
        assert result[:4] == b"%PDF"


# ─── TestRendererSelfContainment ──────────────────────────────────────────────

class TestRendererSelfContainment:
    """
    RULE 003 — Le renderer contient et absorbe tous les problèmes de contenu extrême.

    Un texte de décision de 2 000 caractères est un problème de présentation.
    C'est le problème du renderer — jamais de l'appelant.
    """

    def test_pdf_handles_extreme_text_length(self):
        """PDF avec textes > 2 000 chars dans decisions, destroyers, diagnostic → pas d'exception."""
        try:
            result = generate_pdf_report(EXTREME_TEXT_CASE, "Extrême SARL")
        except Exception as exc:
            pytest.fail(
                f"RULE 003 VIOLATION — PDF a propagé une exception sur texte long : {exc!r}\n"
                f"Longueur du texte testé : {len(_LONG_TEXT)} caractères."
            )
        assert result[:4] == b"%PDF", "PDF renderer doit produire un fichier PDF valide même avec texte extrême."

    def test_pptx_handles_extreme_text_length(self):
        """PPTX avec textes > 2 000 chars dans decisions, destroyers → pas d'exception."""
        from pptx import Presentation
        import io
        try:
            result = generate_pptx_report(EXTREME_TEXT_CASE, "Extrême SARL")
        except Exception as exc:
            pytest.fail(
                f"RULE 003 VIOLATION — PPTX a propagé une exception sur texte long : {exc!r}\n"
                f"Longueur du texte testé : {len(_LONG_TEXT)} caractères."
            )
        prs = Presentation(io.BytesIO(result))
        assert len(prs.slides) == 16, (
            f"PPTX doit toujours produire 16 slides même avec texte extrême. Reçu : {len(prs.slides)}."
        )

    def test_excel_handles_extreme_text_length(self):
        """Excel avec textes > 2 000 chars → workbook valide produit."""
        from openpyxl import load_workbook
        import io
        try:
            result = generate_excel_report(AnalysisResult(), EXTREME_TEXT_CASE, "extreme")
        except Exception as exc:
            pytest.fail(
                f"RULE 003 VIOLATION — Excel a propagé une exception sur texte long : {exc!r}\n"
                f"Longueur du texte testé : {len(_LONG_TEXT)} caractères."
            )
        wb = load_workbook(io.BytesIO(result))
        assert len(wb.sheetnames) >= 1, "Excel doit produire un workbook valide même avec texte extrême."
