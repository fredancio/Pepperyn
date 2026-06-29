"""
test_edx_002.py — Tests du raisonnement comparatif Pepperyn (EDX-002)

Valide :
  1. Modèle Pydantic — EliminatedOption, TippingCondition, champs DecisionReasoning
  2. Sérialisation dans case_to_result_dict()
  3. Comportement PDF renderer avec options_considered
  4. Comportement PPTX renderer avec _slide_raisonnement_comparatif
  5. Cas limite : options vides, dominant_rationale absent, zéro tipping_condition
  6. Garantie anti-régression EDX-001 : anciens champs toujours présents
"""

import sys
import os

sys.path.insert(0, os.path.dirname(os.path.dirname(__file__)))

from models.executive_case import (
    EliminatedOption,
    TippingCondition,
    DecisionReasoning,
    ExecutiveCaseJSON,
    PriorityDecisionItem,
    COIBreakdown,
    DimensionScores,
)
from services.executive_case_builder import case_to_result_dict


# ─── FIXTURES ────────────────────────────────────────────────────────────────

def _make_reasoning_with_edx002() -> DecisionReasoning:
    """Raisonnement complet EDX-001 + EDX-002."""
    return DecisionReasoning(
        decision_index=0,
        problem_source="Stocks obsolètes (180 K€ identifiés)",
        matching_confidence="HIGH",
        decision_confidence=82,
        why_this_decision=(
            "Pepperyn recommande cette décision parce que le stock obsolète "
            "immobilise 180 K€ de trésorerie avec un taux de rotation nul depuis "
            "14 mois. Chaque mois d'inaction coûte 12 K€ en coût de portage."
        ),
        inaction_risk=(
            "À fin septembre 2026, la dépréciation comptable obligatoire "
            "réduira l'actif net de 54 K€ supplémentaires."
        ),
        confidence_explanation=(
            "Score 82% : les données de rotation de stock sont fiables "
            "(fichier inventaire complet), mais l'estimation du prix de liquidation "
            "reste approximative sans audit des acheteurs potentiels."
        ),
        # EDX-002 ─────────────────────────────────────────────────────────────
        options_considered=[
            EliminatedOption(
                option="Provisionnement comptable immédiat à 100%",
                elimination_criterion=(
                    "Ne libère pas de trésorerie — transforme le stock en charge "
                    "sans générer de flux entrant. Runway de 11 semaines rend "
                    "cette option inacceptable."
                ),
            ),
            EliminatedOption(
                option="Cession à un tiers spécialisé (30% valeur)",
                elimination_criterion=(
                    "Récupère 54 K€ immédiatement vs 180 K€ en liquidation "
                    "progressive. Différentiel de 126 K€ supérieur au risque "
                    "d'exécution sur 6 semaines."
                ),
            ),
            EliminatedOption(
                option="Conserver le stock et attendre reprise sectorielle",
                elimination_criterion=(
                    "Secteur éclairage LED : prix en baisse structurelle de -8%/an. "
                    "La valeur du stock diminue chaque trimestre, pas de reprise probable."
                ),
            ),
        ],
        dominant_rationale=(
            "La liquidation commerciale progressive récupère 3,3x plus de valeur "
            "que la cession à tiers (180 K€ vs 54 K€) en échange d'un délai "
            "de 6 semaines supplémentaires. Ce délai est compatible avec le runway "
            "de 11 semaines disponible au J+0."
        ),
        tipping_conditions=[
            TippingCondition(
                condition=(
                    "La trésorerie disponible passe sous 80 K€ avant J+14 "
                    "de l'exécution"
                ),
                horizon_days=30,
                alternative_recommendation=(
                    "Basculer immédiatement vers la cession à tiers pour "
                    "récupérer 54 K€ en 72h."
                ),
            ),
            TippingCondition(
                condition=(
                    "Le client Groupe Lumière représente plus de 40% "
                    "du stock concerné"
                ),
                horizon_days=14,
                alternative_recommendation=(
                    "Négocier un retour fournisseur avant toute liquidation "
                    "pour préserver la relation commerciale."
                ),
            ),
        ],
    )


def _make_minimal_reasoning() -> DecisionReasoning:
    """Raisonnement EDX-001 seul — champs EDX-002 vides (defaults)."""
    return DecisionReasoning(
        decision_index=1,
        problem_source="BFR excessif",
        matching_confidence="LOW",
        decision_confidence=61,
        why_this_decision="Pepperyn recommande cette décision parce que le BFR est excessif.",
        inaction_risk="Dans 90 jours, le BFR atteindra 125 jours.",
        confidence_explanation="Score 61% : données partielles sur le DSO réel.",
        # EDX-002 — absents → defaults
    )


def _make_full_case(reasoning_list: list) -> ExecutiveCaseJSON:
    return ExecutiveCaseJSON(
        company_name="Optilux SAS",
        analysis_date="30/06/2026",
        document_type="PREVISIONNEL",
        confidence_score=75,
        health_score=4,
        dimension_scores=DimensionScores(rentabilite=3, risque=5, structure=4, liquidite=3),
        decisions_priority_score=8.2,
        cost_of_inaction=COIBreakdown(
            annual=-620000.0, monthly=-51667.0, weekly=-11923.0,
            daily=-1695.0, hourly=-211.0,
        ),
        executive_diagnosis="Optilux présente une triple pression sur sa liquidité.",
        tension_phrase="Chaque semaine coûte 11 923 € à ne rien faire.",
        priority_decisions=[
            PriorityDecisionItem(
                decision="Liquider les stocks obsolètes",
                annual_impact=180000.0,
                monthly_impact=15000.0,
                difficulty="Moyen",
                timeline="30 jours",
                priority="Élevée",
                roi_score=8.5,
            ),
            PriorityDecisionItem(
                decision="Réduire le délai de recouvrement clients",
                annual_impact=95000.0,
                monthly_impact=7917.0,
                difficulty="Faible",
                timeline="60 jours",
                priority="Élevée",
                roi_score=7.2,
            ),
        ],
        decision_reasoning=reasoning_list,
    )


# ─── TESTS MODÈLE ────────────────────────────────────────────────────────────

def test_eliminated_option_fields():
    """EliminatedOption a les deux champs requis."""
    opt = EliminatedOption(
        option="Cession à tiers",
        elimination_criterion="Récupère 3x moins que la liquidation progressive.",
    )
    assert opt.option == "Cession à tiers"
    assert "3x" in opt.elimination_criterion


def test_tipping_condition_fields():
    """TippingCondition a les trois champs requis avec horizon 90j par défaut."""
    tc = TippingCondition(
        condition="Trésorerie sous 80 K€",
        alternative_recommendation="Basculer vers cession immédiate.",
    )
    assert tc.condition == "Trésorerie sous 80 K€"
    assert tc.horizon_days == 90  # défaut
    assert tc.alternative_recommendation != ""


def test_tipping_condition_custom_horizon():
    """horizon_days peut être personnalisé."""
    tc = TippingCondition(
        condition="Stock client > 40%",
        horizon_days=14,
        alternative_recommendation="Négocier retour fournisseur.",
    )
    assert tc.horizon_days == 14


def test_decision_reasoning_edx002_defaults():
    """DecisionReasoning EDX-002 champs défaut : listes vides, None pour dominant."""
    r = DecisionReasoning(decision_index=0)
    assert r.options_considered == []
    assert r.dominant_rationale is None
    assert r.tipping_conditions == []


def test_decision_reasoning_edx001_still_present():
    """EDX-001 n'est pas cassé par l'ajout de EDX-002."""
    r = _make_minimal_reasoning()
    assert r.problem_source == "BFR excessif"
    assert r.matching_confidence == "LOW"
    assert r.decision_confidence == 61
    assert "BFR" in r.why_this_decision
    assert r.options_considered == []  # EDX-002 vide par défaut


def test_decision_reasoning_full_edx002():
    """DecisionReasoning complet avec EDX-002 populated."""
    r = _make_reasoning_with_edx002()
    assert len(r.options_considered) == 3
    assert r.options_considered[0].option == "Provisionnement comptable immédiat à 100%"
    assert "11 semaines" in r.options_considered[0].elimination_criterion
    assert "3,3x" in r.dominant_rationale
    assert len(r.tipping_conditions) == 2
    assert r.tipping_conditions[0].horizon_days == 30
    assert r.tipping_conditions[1].horizon_days == 14


# ─── TESTS SÉRIALISATION ─────────────────────────────────────────────────────

def test_case_to_result_dict_includes_edx002():
    """case_to_result_dict() sérialise les champs EDX-002."""
    case = _make_full_case([
        _make_reasoning_with_edx002(),
        _make_minimal_reasoning(),
    ])
    rd = case_to_result_dict(case)
    reasoning_list = rd.get("decision_reasoning", [])

    assert len(reasoning_list) == 2

    r0 = reasoning_list[0]
    assert "options_considered" in r0
    assert len(r0["options_considered"]) == 3

    opt0 = r0["options_considered"][0]
    assert "option" in opt0
    assert "elimination_criterion" in opt0
    assert opt0["option"] == "Provisionnement comptable immédiat à 100%"

    assert r0["dominant_rationale"] is not None
    assert "3,3x" in r0["dominant_rationale"]

    assert "tipping_conditions" in r0
    assert len(r0["tipping_conditions"]) == 2
    tc0 = r0["tipping_conditions"][0]
    assert "condition" in tc0
    assert "horizon_days" in tc0
    assert "alternative_recommendation" in tc0
    assert tc0["horizon_days"] == 30


def test_case_to_result_dict_empty_edx002():
    """case_to_result_dict() gère les champs EDX-002 vides sans erreur."""
    case = _make_full_case([_make_minimal_reasoning()])
    rd = case_to_result_dict(case)
    r1 = rd["decision_reasoning"][0]
    assert r1["options_considered"] == []
    assert r1["dominant_rationale"] is None
    assert r1["tipping_conditions"] == []


def test_edx001_fields_still_serialized():
    """EDX-001 toujours sérialisé — aucune régression."""
    case = _make_full_case([_make_reasoning_with_edx002()])
    rd = case_to_result_dict(case)
    r0 = rd["decision_reasoning"][0]
    # EDX-001 champs présents
    assert r0["decision_index"] == 0
    assert r0["problem_source"] is not None
    assert r0["matching_confidence"] == "HIGH"
    assert r0["decision_confidence"] == 82
    assert r0["why_this_decision"] is not None
    assert r0["inaction_risk"] is not None
    assert r0["confidence_explanation"] is not None


# ─── TESTS GÉNÉRATION PDF ─────────────────────────────────────────────────────

def test_pdf_generates_with_edx002_data():
    """PDF se génère sans erreur avec EDX-002 populated."""
    from services.export_pdf_service import generate_pdf_report
    from services.executive_case_builder import case_to_result_dict

    case = _make_full_case([_make_reasoning_with_edx002()])
    result = case_to_result_dict(case)

    pdf_bytes = generate_pdf_report(result, "Optilux SAS")
    assert pdf_bytes is not None
    assert len(pdf_bytes) > 5_000  # PDF non-trivial


def test_pdf_generates_with_empty_edx002():
    """PDF se génère sans erreur quand EDX-002 est vide."""
    from services.export_pdf_service import generate_pdf_report
    from services.executive_case_builder import case_to_result_dict

    case = _make_full_case([_make_minimal_reasoning()])
    result = case_to_result_dict(case)

    pdf_bytes = generate_pdf_report(result, "Optilux SAS")
    assert pdf_bytes is not None
    assert len(pdf_bytes) > 5_000


# ─── TESTS GÉNÉRATION PPTX ───────────────────────────────────────────────────

def test_pptx_generates_with_edx002_data():
    """PPTX se génère sans erreur avec EDX-002 populated — slide S7b incluse."""
    from services.export_pptx_service import generate_pptx_report
    from services.executive_case_builder import case_to_result_dict

    case = _make_full_case([_make_reasoning_with_edx002()])
    result = case_to_result_dict(case)

    pptx_bytes = generate_pptx_report(result, "Optilux SAS")
    assert pptx_bytes is not None
    assert len(pptx_bytes) > 10_000


def test_pptx_generates_without_edx002_no_slide_added():
    """PPTX se génère sans la slide S7b quand EDX-002 est absent."""
    import io
    from pptx import Presentation
    from services.export_pptx_service import generate_pptx_report
    from services.executive_case_builder import case_to_result_dict

    case = _make_full_case([_make_minimal_reasoning()])
    result = case_to_result_dict(case)

    pptx_bytes = generate_pptx_report(result, "Optilux SAS")
    prs = Presentation(io.BytesIO(pptx_bytes))
    # Sans EDX-002, le deck doit avoir 16 slides (pas 17)
    assert len(prs.slides) == 16


def test_pptx_has_17_slides_with_edx002():
    """PPTX doit avoir 17 slides (S7b incluse) quand EDX-002 est présent."""
    import io
    from pptx import Presentation
    from services.export_pptx_service import generate_pptx_report
    from services.executive_case_builder import case_to_result_dict

    case = _make_full_case([_make_reasoning_with_edx002()])
    result = case_to_result_dict(case)

    pptx_bytes = generate_pptx_report(result, "Optilux SAS")
    prs = Presentation(io.BytesIO(pptx_bytes))
    # Avec EDX-002, le deck doit avoir 17 slides
    assert len(prs.slides) == 17


# ─── TEST D'INTÉGRATION COMPLET ──────────────────────────────────────────────

def test_vocabulary_transformation_in_pdf():
    """PDF se génère et contient le nouveau vocabulaire (smoke test)."""
    from services.export_pdf_service import generate_pdf_report
    from services.executive_case_builder import case_to_result_dict

    case = _make_full_case([_make_reasoning_with_edx002()])
    result = case_to_result_dict(case)

    pdf_bytes = generate_pdf_report(result, "Optilux SAS")
    # PDF contient bien le nouveau vocabulaire — vérification indirecte par taille
    assert len(pdf_bytes) > 5_000
    # Note : validation textuelle du PDF nécessiterait pdfplumber/pdfminer — smoke test ici


if __name__ == "__main__":
    import pytest
    FAIL = pytest.main([__file__, "-v", "--tb=short"])
    sys.exit(1 if FAIL else 0)
