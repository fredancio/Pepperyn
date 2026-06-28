"""
tests/test_rule_002_zero_truncation.py

PEPPERYN PRODUCT QUALITY CONTRACT — RULE 002
Zero Truncated Content

Ce test est le gardien permanent de RULE 002.
Le critère n'est pas "le texte est moins tronqué".
Le critère est : aucune information n'est sacrifiée pour des raisons de mise en page.

Protocole de vérification :
  1. Code source   — aucun slice de contenu texte [:n] dans les renderers
  2. PPTX décisions — chaque décision de l'ExecutiveCase est intégralement présente
  3. PPTX risques   — chaque risque est intégralement présent
  4. PPTX destroyers — chaque destructeur est intégralement présent
  5. PDF décisions  — chaque décision est intégralement présente
  6. PDF destroyers — chaque destructeur est intégralement présent
  7. Tables python-pptx — row heights calculées dynamiquement (>= minimum)
  8. Invariant structurel — les helpers _auto_row_h et _fit_table_rows existent

RULE 004 — Renderers DISPLAY. They never THINK.
RULE 002 — Zero Truncated Content.
"""
import sys, os, re, subprocess, pytest
sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))
os.environ.setdefault("ANTHROPIC_API_KEY", "dummy")

from pptx import Presentation
from models.schemas import AnalysisResult
from models.executive_case import (
    ExecutiveCaseJSON, DimensionScores, COIBreakdown, KPICard, ValueDestroyerItem,
    PriorityDecisionItem, ExecutionLogItem, ProjectionSeries, Scenarios, ScenarioItem,
    RiskItem, DataQuality,
)
from services.excel_export import generate_excel_report
from services.export_pdf_service import generate_pdf_report
from services.export_pptx_service import generate_pptx_report, _auto_row_h, _fit_table_rows

# ─── Cas Optilux SAS (décisions volontairement longues pour tester les limites) ─
OPTILUX = ExecutiveCaseJSON(
    company_name="Optilux SAS", analysis_date="28/06/2026", document_type="PREVISIONNEL",
    confidence_score=81, health_score=3, decisions_priority_score=8.4,
    dimension_scores=DimensionScores(rentabilite=2, risque=3, structure=4, liquidite=2),
    cost_of_inaction=COIBreakdown(annual=-1_699_000, monthly=-141_583, weekly=-32_673, daily=-4_654, hourly=-581),
    executive_diagnosis="Optilux SAS traverse une crise de rentabilité structurelle aggravée par une tension de trésorerie imminente. Avec seulement 180 K€ de trésorerie disponible et un EBITDA négatif de -240 K€, l'entreprise dispose de moins de 11 semaines de runway à charge constante.",
    tension_phrase="Chaque mois d'inaction coûte 141 583 € et rapproche Optilux d'un point de non-retour.",
    inaction_risk="À trésorerie et trajectoire constantes, Optilux atteindra l'insolvabilité technique dans 11 semaines.",
    structural_loss_statement="La perte structurelle annuelle d'Optilux est estimée à 1,7 M€, soit 20,7 % du CA.",
    structural_loss_value=-1_699_000, urgency_level="CRITIQUE — Trésorerie à 11 semaines",
    value_creation_statement="Si les quatre décisions prioritaires sont engagées dans les 30 à 90 prochains jours, Optilux peut récupérer entre 820 K€ et 1,15 M€ d'impact annualisé.",
    kpi_dashboard=[
        KPICard(label="Chiffre d'affaires", value="8,2 M€"),
        KPICard(label="EBITDA", value="-240 K€", status="missing"),
        KPICard(label="Trésorerie", value="180 K€", status="missing"),
        KPICard(label="Marge brute", value="28 %", status="missing"),
        KPICard(label="DSO clients", value="87 jours", status="missing"),
        KPICard(label="Stock obsolète", value="620 K€", status="missing"),
    ],
    value_destroyers=[
        ValueDestroyerItem(name="Stock obsolète non provisionné valorisé à 620 K€ (3 ans d'accumulation)", annual_impact=-620_000, monthly_impact=-51_667, trend="↑"),
        ValueDestroyerItem(name="Érosion des marges commerciales (-7 pts de marge brute en 18 mois)", annual_impact=-574_000, monthly_impact=-47_833, trend="↑"),
        ValueDestroyerItem(name="Sous-performance commerciale par rapport au plan initial (-18 % vs budget)", annual_impact=-320_000, monthly_impact=-26_667, trend="→"),
        ValueDestroyerItem(name="Surcoût BFR dû à un DSO excessif (87 jours vs 45 jours standard sectoriel)", annual_impact=-185_000, monthly_impact=-15_417, trend="↑"),
    ],
    priority_decisions=[
        PriorityDecisionItem(
            decision="Provisionner et liquider le stock obsolète (620 K€) avant la clôture de septembre — vente flash à -40 % auprès d'installateurs partenaires",
            annual_impact=400_000, monthly_impact=33_333, difficulty="Faible", timeline="30 jours", priority="Élevée", roi_score=9.2, owner="DAF + Directeur Logistique", status="À lancer"),
        PriorityDecisionItem(
            decision="Plan de recouvrement agressif — réduire le DSO de 87j à 55j en 90 jours via relances structurées et remises d'escompte 2/10 net 45",
            annual_impact=185_000, monthly_impact=15_417, difficulty="Faible", timeline="30 jours", priority="Élevée", roi_score=8.7, owner="DAF", status="À lancer"),
        PriorityDecisionItem(
            decision="Renégocier les conditions tarifaires avec les 3 fournisseurs principaux (70 % des achats) — objectif +3 pts de marge brute via volumes garantis",
            annual_impact=246_000, monthly_impact=20_500, difficulty="Moyen", timeline="60 jours", priority="Élevée", roi_score=7.8, owner="Directeur Commercial", status="À lancer"),
        PriorityDecisionItem(
            decision="Restructurer le plan de commissionnement commercial — objectifs trimestriels + bonus collectif sur marge (pas sur CA)",
            annual_impact=274_000, monthly_impact=22_833, difficulty="Moyen", timeline="90 jours", priority="Moyenne", roi_score=6.4, owner="DG + RH", status="À lancer"),
    ],
    roadmap_30_60_90={
        "30": ["Audit complet du stock", "Relances clients prioritaires"],
        "60": ["Finaliser liquidation stock obsolète"],
        "90": ["Déployer nouvelle grille commissionnement"],
    },
    execution_log=[
        ExecutionLogItem(decision="Audit stock obsolète", owner="Directeur Logistique", impact=400_000, due_date="28/07/2026", difficulty="Faible", roi_score=9.2, status="À lancer"),
    ],
    series=ProjectionSeries(
        action=[33_333,66_666,155_000,195_000,235_000,285_000,325_000,365_000,420_000,480_000,545_000,620_000],
        inaction=[-141_583,-283_166,-424_749,-566_332,-707_915,-849_498,-991_081,-1_132_664,-1_274_247,-1_415_830,-1_557_413,-1_699_000],
        equilibrium=[-108_250,-68_000,45_000,115_000,195_000,285_000,340_000,395_000,450_000,505_000,560_000,615_000],
    ),
    scenarios=Scenarios(
        best=ScenarioItem(label="Exécution rapide (J+30)", description="Les 4 décisions engagées en 30j. Impact annualisé : +820 K€."),
        likely=ScenarioItem(label="Exécution standard (J+60)", description="Décisions stock+DSO en J+30, fournisseurs en J+60."),
        worst=ScenarioItem(label="Exécution partielle", description="Seule la liquidation du stock est engagée (+372 K€ cash court terme)."),
    ),
    major_risks=[
        RiskItem(
            description="Rupture de trésorerie avant fin T3 2026 si liquidation stock et plan recouvrement ne génèrent pas 300 K€ minimum d'ici fin août",
            severity="Élevé", impact="Critique", horizon="Immédiat (< 90 jours)"),
        RiskItem(
            description="Perte du client Groupe Lumière (24 % du CA, 1,97 M€) suite à défaillance de livraison provoquée par tension de trésorerie",
            severity="Élevé", impact="Fort", horizon="Court terme (3-6 mois)"),
    ],
    data_quality=DataQuality(score=81, anomalies=["Données Q2 2026 partielles (clôture en cours)"], assumptions=["Taux de réussite liquidation stock estimé à 60 %"], limits=["Filiales hors périmètre"]),
)


@pytest.fixture(scope="module")
def generated_files(tmp_path_factory):
    d = tmp_path_factory.mktemp("rule002")
    pdf_path  = d / "Optilux.pdf"
    pptx_path = d / "Optilux.pptx"
    xl_path   = d / "Optilux.xlsx"
    pdf_path.write_bytes(generate_pdf_report(OPTILUX, "Optilux SAS"))
    pptx_path.write_bytes(generate_pptx_report(OPTILUX, "Optilux SAS"))
    xl_path.write_bytes(generate_excel_report(AnalysisResult(), OPTILUX, "Optilux"))
    return {"pdf": pdf_path, "pptx": pptx_path, "xl": xl_path}


CONTENT_SLICE_PATTERN = re.compile(
    r'\.(decision|description|name|label|value|text|diagnostic|risque|desc|diag|act'
    r'|dec_txt|dec_text|val_str|a|h)\[:\d+\]'
)


class TestNoContentSlicesInSource:
    """RULE 002 — Le code source ne doit contenir aucun slice de contenu textuel."""

    RENDERER_FILES = [
        "services/export_pptx_service.py",
        "services/export_pdf_service.py",
    ]

    def test_no_content_slices_in_pptx_renderer(self):
        base = os.path.join(os.path.dirname(__file__), "..")
        fpath = os.path.join(base, "services/export_pptx_service.py")
        violations = self._find_violations(fpath)
        assert not violations, (
            f"RULE 002 VIOLATION — export_pptx_service.py contient des slices de contenu :\n"
            + "\n".join(f"  L{n}: {c}" for n, c in violations)
        )

    def test_no_content_slices_in_pdf_renderer(self):
        base = os.path.join(os.path.dirname(__file__), "..")
        fpath = os.path.join(base, "services/export_pdf_service.py")
        violations = self._find_violations(fpath)
        assert not violations, (
            f"RULE 002 VIOLATION — export_pdf_service.py contient des slices de contenu :\n"
            + "\n".join(f"  L{n}: {c}" for n, c in violations)
        )

    def _find_violations(self, fpath: str):
        violations = []
        for i, line in enumerate(open(fpath), 1):
            if line.strip().startswith("#"):
                continue
            if CONTENT_SLICE_PATTERN.search(line):
                violations.append((i, line.strip()))
        return violations


class TestPPTXCompleteness:
    """RULE 002 — Toutes les informations de l'ExecutiveCase sont intégralement dans le PPTX."""

    def _pptx_all_text(self, pptx_path) -> str:
        prs = Presentation(str(pptx_path))
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    parts.append(shape.text_frame.text.replace("\n", " "))
                elif shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            parts.append(cell.text_frame.text.replace("\n", " "))
        return " ".join(parts)

    def test_all_decisions_complete_in_pptx(self, generated_files):
        text = self._pptx_all_text(generated_files["pptx"])
        text_norm = " ".join(text.split())
        truncated = []
        for i, dec in enumerate(OPTILUX.priority_decisions, 1):
            dec_norm = " ".join(dec.decision.split())
            if dec_norm not in text_norm:
                truncated.append(f"D{i}: {dec.decision[:80]}...")
        assert not truncated, (
            f"RULE 002 VIOLATION — PPTX : décisions tronquées :\n"
            + "\n".join(f"  {d}" for d in truncated)
        )

    def test_all_risks_complete_in_pptx(self, generated_files):
        text = self._pptx_all_text(generated_files["pptx"])
        text_norm = " ".join(text.split())
        truncated = []
        for i, risk in enumerate(OPTILUX.major_risks, 1):
            risk_norm = " ".join(risk.description.split())
            if risk_norm not in text_norm:
                truncated.append(f"R{i}: {risk.description[:80]}...")
        assert not truncated, (
            f"RULE 002 VIOLATION — PPTX : risques tronqués :\n"
            + "\n".join(f"  {r}" for r in truncated)
        )

    def test_all_destroyers_complete_in_pptx(self, generated_files):
        text = self._pptx_all_text(generated_files["pptx"])
        text_norm = " ".join(text.split())
        truncated = []
        for i, d in enumerate(OPTILUX.value_destroyers, 1):
            name_norm = " ".join(d.name.split())
            if name_norm not in text_norm:
                truncated.append(f"VD{i}: {d.name[:80]}...")
        assert not truncated, (
            f"RULE 002 VIOLATION — PPTX : destructeurs tronqués :\n"
            + "\n".join(f"  {d}" for d in truncated)
        )


class TestPDFCompleteness:
    """RULE 002 — Toutes les informations de l'ExecutiveCase sont intégralement dans le PDF."""

    def _pdf_text(self, pdf_path) -> str:
        result = subprocess.run(["pdftotext", str(pdf_path), "-"],
                                capture_output=True, text=True)
        return " ".join(result.stdout.split())

    def test_all_decisions_complete_in_pdf(self, generated_files):
        text = self._pdf_text(generated_files["pdf"])
        truncated = []
        for i, dec in enumerate(OPTILUX.priority_decisions, 1):
            dec_norm = " ".join(dec.decision.split())
            if dec_norm not in text:
                truncated.append(f"D{i}: {dec.decision[:80]}...")
        assert not truncated, (
            f"RULE 002 VIOLATION — PDF : décisions tronquées :\n"
            + "\n".join(f"  {d}" for d in truncated)
        )

    def test_all_destroyers_complete_in_pdf(self, generated_files):
        text = self._pdf_text(generated_files["pdf"])
        truncated = []
        for i, d in enumerate(OPTILUX.value_destroyers, 1):
            name_norm = " ".join(d.name.split())
            if name_norm not in text:
                truncated.append(f"VD{i}: {d.name[:80]}...")
        assert not truncated, (
            f"RULE 002 VIOLATION — PDF : destructeurs tronqués :\n"
            + "\n".join(f"  {d}" for d in truncated)
        )


class TestAutoRowHeightHelper:
    """RULE 002 — _auto_row_h() calcule des hauteurs suffisantes pour le contenu."""

    def test_short_text_gets_minimum_height(self):
        # 10 chars dans une colonne de 3" à 10pt → 1 ligne → hauteur min
        col_w = int(3 * 914400)  # 3 inches en EMU
        h = _auto_row_h("Court", col_w, 10, base_pt=28.0)
        assert h >= int(Pt(28)), "Hauteur minimale non respectée"

    def test_long_decision_gets_adequate_height(self):
        from pptx.util import Pt
        decision = "Provisionner et liquider le stock obsolète (620 K€) avant la clôture de septembre — vente flash à -40 % auprès d'installateurs partenaires"
        col_w = int(3.2 * 914400)  # 3.2 inches (26% of 12.33" slide)
        h = _auto_row_h(decision, col_w, 10)
        # 140 chars en 3.2" à 10pt → >1 ligne → hauteur > minimum
        assert h > int(Pt(28)), f"Hauteur insuffisante pour décision longue : {h} EMU"

    def test_very_long_text_gets_proportional_height(self):
        from pptx.util import Pt
        text = "A" * 200
        col_w = int(3.0 * 914400)
        h1 = _auto_row_h(text[:50],  col_w, 10)
        h2 = _auto_row_h(text[:200], col_w, 10)
        assert h2 > h1, "Les textes plus longs doivent produire des hauteurs plus grandes"


try:
    from pptx.util import Pt
except ImportError:
    pass
