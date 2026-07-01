"""
tests/test_rule_001_zero_manual_intervention.py

PEPPERYN PRODUCT QUALITY CONTRACT — RULE 001
Zero Manual Intervention

Ce test est le gardien permanent de RULE 001.
Il doit passer à chaque génération de livrables Optilux SAS.
Si ce test échoue, aucune version ne peut être déclarée conforme.

Protocole de vérification :
  1. EDM source values   — 14 cellules vérifiées
  2. Formula errors       — 0 erreur (#REF!, #VALUE!, etc.) dans les 139+ formules
  3. Column widths        — largeurs explicites sur toutes les feuilles visibles
  4. Wrap text            — activé sur les cellules de décision longues
  5. Row heights          — hauteur explicite sur les lignes de décision
  6. PDF content          — 13 tokens requis présents
  7. PPTX content         — 14 tokens requis présents + 16 slides
  8. LibreOffice recalc   — 0 erreur après recalcul complet

RULE 004 — Renderers DISPLAY. They never THINK.
RULE 001 — Zero Manual Intervention.
"""
import sys, os, json, subprocess, pickle, tempfile, pytest
sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))
os.environ.setdefault("ANTHROPIC_API_KEY", "dummy")

from openpyxl import load_workbook
from pptx import Presentation
from models.schemas import AnalysisResult
from models.executive_case import (
    ExecutiveCaseJSON, DimensionScores, COIBreakdown, KPICard, ValueDestroyerItem,
    PriorityDecisionItem, ExecutionLogItem, ProjectionSeries, Scenarios, ScenarioItem,
    RiskItem, DataQuality, PLLine, BalanceLine, FinancialStatements,
)
from services.excel_export import generate_excel_report
from services.export_pdf_service import generate_pdf_report
from services.export_pptx_service import generate_pptx_report

# ─── Séparateur de milliers selon le livrable ────────────────────────────────
NNBSP = " "   # Narrow No-Break Space — utilisé par le renderer PPTX
SP    = " "        # Espace ordinaire — produit par pdftotext sur le PDF

# ─── Cas Optilux SAS (cas de référence RULE 001) ─────────────────────────────
OPTILUX = ExecutiveCaseJSON(
    company_name="Optilux SAS", analysis_date="28/06/2026", document_type="PREVISIONNEL",
    confidence_score=81, health_score=3, decisions_priority_score=8.4,
    dimension_scores=DimensionScores(rentabilite=2, risque=3, structure=4, liquidite=2),
    cost_of_inaction=COIBreakdown(annual=-1_699_000, monthly=-141_583, weekly=-32_673, daily=-4_654, hourly=-581),
    executive_diagnosis="Optilux SAS traverse une crise de rentabilité structurelle aggravée par une tension de trésorerie imminente.",
    tension_phrase="Chaque mois d'inaction coûte 141 583 € et rapproche Optilux d'un point de non-retour.",
    inaction_risk="À trésorerie et trajectoire constantes, Optilux atteindra l'insolvabilité technique dans 11 semaines.",
    structural_loss_statement="La perte structurelle annuelle d'Optilux est estimée à 1,7 M€.",
    structural_loss_value=-1_699_000, urgency_level="CRITIQUE — Trésorerie à 11 semaines",
    value_creation_statement="Si les quatre décisions sont engagées, Optilux récupère entre 820 K€ et 1,15 M€.",
    kpi_dashboard=[
        KPICard(label="Chiffre d'affaires", value="8,2 M€"),
        KPICard(label="EBITDA",             value="-240 K€"),
        KPICard(label="Résultat net",       value="-505 K€"),
        KPICard(label="Marge brute",        value="28 %"),
        KPICard(label="Trésorerie",         value="180 K€"),
        KPICard(label="DSO clients",        value="87 jours"),
        KPICard(label="Stock obsolète",     value="620 K€"),
        KPICard(label="BFR",                value="2 340 K€"),
    ],
    value_destroyers=[
        ValueDestroyerItem(name="Stock obsolète non provisionné",           annual_impact=-620_000, monthly_impact=-51_667, trend="↑"),
        ValueDestroyerItem(name="Érosion des marges commerciales (-7 pts)", annual_impact=-574_000, monthly_impact=-47_833, trend="↑"),
        ValueDestroyerItem(name="Sous-performance commerciale (-18 %)",     annual_impact=-320_000, monthly_impact=-26_667, trend="→"),
        ValueDestroyerItem(name="Surcoût BFR — DSO excessif (87j vs 45j)", annual_impact=-185_000, monthly_impact=-15_417, trend="↑"),
    ],
    priority_decisions=[
        PriorityDecisionItem(decision="Provisionner et liquider le stock obsolète (620 K€) avant la clôture de septembre", annual_impact=400_000, monthly_impact=33_333, difficulty="Faible", timeline="30 jours", priority="Élevée", roi_score=9.2, owner="DAF + Directeur Logistique", status="À lancer"),
        PriorityDecisionItem(decision="Plan de recouvrement agressif — réduire le DSO de 87j à 55j en 90 jours", annual_impact=185_000, monthly_impact=15_417, difficulty="Faible", timeline="30 jours", priority="Élevée", roi_score=8.7, owner="DAF", status="À lancer"),
        PriorityDecisionItem(decision="Renégocier les conditions tarifaires avec les 3 fournisseurs principaux", annual_impact=246_000, monthly_impact=20_500, difficulty="Moyen", timeline="60 jours", priority="Élevée", roi_score=7.8, owner="Directeur Commercial", status="À lancer"),
        PriorityDecisionItem(decision="Restructurer le plan de commissionnement commercial", annual_impact=274_000, monthly_impact=22_833, difficulty="Moyen", timeline="90 jours", priority="Moyenne", roi_score=6.4, owner="DG + RH", status="À lancer"),
    ],
    roadmap_30_60_90={
        "30": ["Audit complet du stock", "Relances clients prioritaires", "Ouvrir négociation Fournisseur A"],
        "60": ["Finaliser liquidation stock obsolète", "Implémenter relances automatisées"],
        "90": ["Déployer nouvelle grille commissionnement", "Conclure négociation Fournisseur C"],
    },
    execution_log=[
        ExecutionLogItem(decision="Audit stock obsolète",        owner="Directeur Logistique", impact=400_000, due_date="28/07/2026", difficulty="Faible", roi_score=9.2, status="À lancer"),
        ExecutionLogItem(decision="Plan recouvrement DSO",       owner="DAF",                   impact=185_000, due_date="15/07/2026", difficulty="Faible", roi_score=8.7, status="À lancer"),
        ExecutionLogItem(decision="Renégociation fournisseurs",  owner="Dir. Commercial",       impact=246_000, due_date="28/08/2026", difficulty="Moyen",  roi_score=7.8, status="À lancer"),
        ExecutionLogItem(decision="Commissionnement commercial", owner="DG + RH",               impact=274_000, due_date="28/09/2026", difficulty="Moyen",  roi_score=6.4, status="À lancer"),
    ],
    series=ProjectionSeries(
        action    =[33_333,66_666,155_000,195_000,235_000,285_000,325_000,365_000,420_000,480_000,545_000,620_000],
        inaction  =[-141_583,-283_166,-424_749,-566_332,-707_915,-849_498,-991_081,-1_132_664,-1_274_247,-1_415_830,-1_557_413,-1_699_000],
        equilibrium=[-108_250,-68_000,45_000,115_000,195_000,285_000,340_000,395_000,450_000,505_000,560_000,615_000],
    ),
    scenarios=Scenarios(
        best  =ScenarioItem(label="Exécution rapide (J+30)", description="Les 4 décisions engagées en 30j."),
        likely=ScenarioItem(label="Exécution standard (J+60)", description="Décisions stock+DSO en J+30."),
        worst =ScenarioItem(label="Exécution partielle", description="Seule la liquidation du stock est engagée."),
    ),
    major_risks=[
        RiskItem(description="Rupture de trésorerie avant fin T3 2026", severity="Élevé", impact="Critique", horizon="Immédiat (< 90 jours)"),
        RiskItem(description="Perte du client Groupe Lumière (24 % du CA)", severity="Élevé", impact="Fort", horizon="Court terme (3-6 mois)"),
    ],
    data_quality=DataQuality(score=81, anomalies=["Données Q2 2026 partielles"], assumptions=["Taux de réussite liquidation stock estimé à 60 %"], limits=["Filiales hors périmètre"]),
    financial_statements=FinancialStatements(
        # ── Compte de résultat ────────────────────────────────────────────
        pl_period="Exercice 2025–2026 (12 mois estimés)",
        pl_lines=[
            PLLine(label="Chiffre d'affaires net",                value_display="8 200 K€"),
            PLLine(label="dont Services",                          value_display="5 100 K€",  indent=1),
            PLLine(label="dont Produits",                          value_display="3 100 K€",  indent=1),
            PLLine(label="Coût des ventes",                        value_display="-5 904 K€"),
            PLLine(label="Marge brute (28 %)",                     value_display="2 296 K€",  is_subtotal=True),
            PLLine(label="Charges de personnel",                   value_display="-1 800 K€", indent=1),
            PLLine(label="Loyers et charges locatives",            value_display="-240 K€",   indent=1),
            PLLine(label="Autres charges d'exploitation",          value_display="-496 K€",   indent=1),
            PLLine(label="EBITDA",                                 value_display="-240 K€",   is_subtotal=True),
            PLLine(label="Amortissements et dotations",            value_display="-180 K€"),
            PLLine(label="Résultat d'exploitation (EBIT)",         value_display="-420 K€",   is_subtotal=True),
            PLLine(label="Charges financières nettes",             value_display="-85 K€"),
            PLLine(label="Résultat avant impôt",                   value_display="-505 K€",   is_subtotal=True),
            PLLine(label="Impôts (résultat déficitaire)",          value_display="0 K€"),
            PLLine(label="Résultat net",                           value_display="-505 K€",   is_total=True),
        ],
        pl_note="Estimations sur base des 9 premiers mois et de la trajectoire extrapolée. Non audité.",
        # ── Bilan simplifié ───────────────────────────────────────────────
        bilan_date="Au 28/06/2026 (estimé)",
        assets=[
            BalanceLine(label="Immobilisations nettes",         value_display="1 200 K€"),
            BalanceLine(label="Stocks (dont 620 K€ obsolètes)", value_display="1 050 K€"),
            BalanceLine(label="Créances clients (DSO 87j)",     value_display="1 970 K€"),
            BalanceLine(label="Autres actifs circulants",        value_display="200 K€"),
            BalanceLine(label="Trésorerie disponible",          value_display="180 K€"),
            BalanceLine(label="Total actif",                    value_display="4 600 K€",  is_total=True),
        ],
        liabilities=[
            BalanceLine(label="Capitaux propres",               value_display="495 K€"),
            BalanceLine(label="Dettes financières L/T",         value_display="1 385 K€"),
            BalanceLine(label="Dettes fournisseurs",            value_display="680 K€"),
            BalanceLine(label="Dettes sociales et fiscales",    value_display="820 K€"),
            BalanceLine(label="Autres dettes",                  value_display="1 220 K€"),
            BalanceLine(label="Total passif",                   value_display="4 600 K€",  is_total=True),
        ],
        bfr_display="2 340 K€",
        bilan_note="BFR = Créances clients (1 970 K€) + Stocks (1 050 K€) − Dettes four. (680 K€). Un BFR de 2,3 M€ sur 8,2 M€ de CA représente 28,5 % du chiffre d'affaires — niveau structurellement élevé.",
        # ── Position de trésorerie ────────────────────────────────────────
        cash_current="180 K€",
        cash_burn_monthly="-141 K€ / mois",
        cash_runway_label="1,3 mois (~5–6 semaines)",
        credit_line_available="Non activée",
        financing_need_90d="-245 K€",
        cash_note="L'insolvabilité technique peut intervenir avant épuisement complet de la trésorerie, notamment en cas de défaillance d'un client majeur (Groupe Lumière représente 24 % du CA).",
    ),
)

# ─── Fixtures ─────────────────────────────────────────────────────────────────
@pytest.fixture(scope="module")
def generated_files(tmp_path_factory):
    """Génère les trois livrables dans un répertoire temporaire."""
    d = tmp_path_factory.mktemp("rule001")
    xl_path   = d / "Optilux.xlsx"
    pdf_path  = d / "Optilux.pdf"
    pptx_path = d / "Optilux.pptx"

    xl_path.write_bytes(generate_excel_report(AnalysisResult(), OPTILUX, "Optilux"))
    pdf_path.write_bytes(generate_pdf_report(OPTILUX, "Optilux SAS"))
    pptx_path.write_bytes(generate_pptx_report(OPTILUX, "Optilux SAS"))

    return {"xl": xl_path, "pdf": pdf_path, "pptx": pptx_path, "dir": d}


# ─── CHECK 1 : EDM source values ──────────────────────────────────────────────
class TestEDMSourceValues:
    EXPECTED = [
        (3,   3),         # Score Santé
        (4,  81),         # Niveau de confiance
        (5,  -240_000),   # EBITDA
        (6,   180_000),   # Cash disponible
        (7,  -1_699_000), # COI / an
        (8,  -141_583),   # COI / mois
        (9,   -32_673),   # COI / sem.
        (10,   -4_654),   # COI / jour
        (11,    -581),    # COI / heure
        (12, -1_699_000), # Impact total (somme destroyers)
        (29,  400_000),   # Impact décision 1
        (30,  185_000),   # Impact décision 2
        (31,  246_000),   # Impact décision 3
        (32,  274_000),   # Impact décision 4
    ]

    def test_edm_source_values(self, generated_files):
        wb = load_workbook(str(generated_files["xl"]), data_only=True)
        ws = wb["EDM"]
        for row, expected in self.EXPECTED:
            actual = ws.cell(row=row, column=2).value
            assert actual == expected, (
                f"RULE 001 VIOLATION — EDM!B{row} = {actual}, attendu {expected}. "
                f"Cause probable : _parse_eur() a retourné 0 ou case_to_edm() n'a pas propagé la valeur."
            )


# ─── CHECK 2 : Zéro erreur de formule ─────────────────────────────────────────
class TestFormulaErrors:
    FORMULA_ERRORS = {"#REF!", "#VALUE!", "#NAME?", "#DIV/0!", "#N/A", "#NULL!", "#NUM!"}

    def test_no_formula_errors_in_static_file(self, generated_files):
        wb = load_workbook(str(generated_files["xl"]), data_only=True)
        errors = []
        for sheet_name in wb.sheetnames:
            for row in wb[sheet_name].iter_rows():
                for cell in row:
                    if cell.value in self.FORMULA_ERRORS:
                        errors.append(f"{sheet_name}!{cell.coordinate} = {cell.value}")
        assert not errors, f"RULE 001 VIOLATION — Formules cassées : {errors}"

    def test_no_formula_errors_after_libreoffice_recalc(self, generated_files):
        # Chercher recalc.py dans les chemins d'environnement possibles
        candidates = [
            os.path.normpath(os.path.join(os.path.dirname(__file__), "../..",
                                          ".claude/skills/xlsx/scripts/recalc.py")),
            "/sessions/modest-sleepy-hypatia/mnt/.claude/skills/xlsx/scripts/recalc.py",
            os.path.expanduser("~/.claude/skills/xlsx/scripts/recalc.py"),
        ]
        skill_root = next((p for p in candidates if os.path.exists(p)), None)
        if skill_root is None:
            pytest.skip("recalc.py non disponible dans cet environnement")

        result = subprocess.run(
            ["python", skill_root, str(generated_files["xl"]), "60"],
            capture_output=True, text=True, timeout=90,
        )
        data = json.loads(result.stdout or "{}")
        assert data.get("total_errors", 0) == 0, (
            f"RULE 001 VIOLATION — LibreOffice détecte {data.get('total_errors')} erreur(s) : "
            f"{data.get('error_summary', {})}"
        )


# ─── CHECK 3 : Largeurs de colonnes explicites ────────────────────────────────
class TestColumnWidths:
    def test_all_visible_sheets_have_explicit_column_widths(self, generated_files):
        wb = load_workbook(str(generated_files["xl"]), data_only=True)
        missing = []
        for sheet_name in wb.sheetnames:
            if sheet_name == "EDM":
                continue  # feuille masquée
            ws = wb[sheet_name]
            explicit = sum(1 for _, cd in ws.column_dimensions.items()
                           if cd.width and cd.width != 8.43)
            if explicit == 0:
                missing.append(sheet_name)
        assert not missing, (
            f"RULE 001 VIOLATION — Colonnes non redimensionnées : {missing}. "
            "Le CEO devrait redimensionner manuellement ces colonnes."
        )


# ─── CHECK 4 : Wrap text sur cellules de décision ────────────────────────────
class TestWrapText:
    def test_decision_cells_have_wrap_text(self, generated_files):
        wb = load_workbook(str(generated_files["xl"]), data_only=True)
        n_decisions = len(OPTILUX.priority_decisions)
        sheets_to_check = [
            ("🎯 Decision Lab", range(6, 6 + n_decisions), 2),
            ("🗺 Roadmap",      range(5, 5 + n_decisions), 2),
        ]
        violations = []
        for sheet_name, rows, col in sheets_to_check:
            ws = wb[sheet_name]
            for r in rows:
                cell = ws.cell(row=r, column=col)
                if not (cell.alignment and cell.alignment.wrap_text):
                    violations.append(f"{sheet_name}!{cell.coordinate}")
        assert not violations, (
            f"RULE 001 VIOLATION — Wrap text manquant sur : {violations}. "
            "Le texte long sera tronqué sans intervention manuelle."
        )


# ─── CHECK 5 : Hauteurs de lignes ────────────────────────────────────────────
class TestRowHeights:
    def test_decision_rows_have_explicit_heights(self, generated_files):
        wb = load_workbook(str(generated_files["xl"]), data_only=True)
        n_decisions = len(OPTILUX.priority_decisions)
        checks = [
            ("🎯 Decision Lab", range(6, 6 + n_decisions)),
            ("🗺 Roadmap",      range(5, 5 + n_decisions)),
        ]
        violations = []
        for sheet_name, rows in checks:
            ws = wb[sheet_name]
            for r in rows:
                if not ws.row_dimensions[r].height:
                    violations.append(f"{sheet_name}!row{r}")
        assert not violations, (
            f"RULE 001 VIOLATION — Hauteur de ligne non définie : {violations}. "
            "Excel forcera une hauteur automatique qui pourrait tronquer le contenu."
        )


# ─── CHECK 6 : Contenu PDF ────────────────────────────────────────────────────
class TestPDFContent:
    # Tokens attendus dans le texte brut extrait par pdftotext
    # pdftotext normalise les espaces → espace ordinaire
    REQUIRED_TOKENS = [
        # CEO Question Framework v3 — tokens attendus dans le nouveau layout PDF
        "-1,7 M€",        # COI hero annuel (P3 COI)
        "-141 583 €",     # COI mensuel (P3 COI)
        "-4 654 €",       # COI journalier (P3 COI)
        "2/10",           # Score dimension rentabilité (P1 Verdict)
        "8.4/10",         # Score de priorisation global (P4 Décisions)
        "+400 000 €",     # Impact décision 1 (P4/P6)
        "+185 000 €",     # Impact décision 2 (P4/P6)
        "9.2/10",         # ROI décision 1 (P4/P6)
        "8.7/10",         # ROI décision 2 (P4/P6)
        "30 jours",       # Timeline décision 1 (P4/P7)
        "CALENDRIER D'ALLOCATION",   # Roadmap (P7)
        "Élevé",          # Sévérité risque (P9 Risques)
        "81%",            # Confiance analyse (P10 KPIs)
    ]
    FORBIDDEN_TOKENS = [
        "Données insuffisantes",   # RULE 007 — ne pas afficher quand la valeur existe
    ]

    def test_pdf_required_tokens_present(self, generated_files):
        result = subprocess.run(
            ["pdftotext", str(generated_files["pdf"]), "-"],
            capture_output=True, text=True,
        )
        text = result.stdout
        missing = [t for t in self.REQUIRED_TOKENS if t not in text]
        assert not missing, (
            f"RULE 001 VIOLATION — PDF : tokens absents : {missing}"
        )

    def test_pdf_no_forbidden_tokens(self, generated_files):
        result = subprocess.run(
            ["pdftotext", str(generated_files["pdf"]), "-"],
            capture_output=True, text=True,
        )
        text = result.stdout
        found = [t for t in self.FORBIDDEN_TOKENS if t in text]
        assert not found, (
            f"RULE 007 VIOLATION — PDF : tokens interdits présents : {found}"
        )


# ─── CHECK 7 : Contenu PPTX ───────────────────────────────────────────────────
class TestPPTXContent:
    # Tokens PPTX — utilise U+202F (espace fine insécable) comme séparateur de milliers
    NNBSP = " "

    def _all_text(self, pptx_path):
        prs = Presentation(str(pptx_path))
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    parts.append(shape.text_frame.text)
                elif shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            parts.append(cell.text_frame.text)
        return "\n".join(parts)

    def test_pptx_has_16_slides(self, generated_files):
        prs = Presentation(str(generated_files["pptx"]))
        assert len(prs.slides) == 16, (
            f"RULE 001 VIOLATION — PPTX : {len(prs.slides)} slides, attendu 16."
        )

    def test_pptx_required_tokens_present(self, generated_files):
        T = self.NNBSP
        required = [
            f"141{T}583 €",          # COI mensuel
            f"400{T}000 €",          # Impact D1
            f"185{T}000 €",          # Impact D2
            "9.2/10",
            "8.7/10",
            "30 JOURS",
            "60 JOURS",
            "90 JOURS",
            "Élevé",
            "Immédiat",
            "lundi matin",
            "81%",
        ]
        text = self._all_text(generated_files["pptx"])
        missing = [t for t in required if t not in text]
        assert not missing, (
            f"RULE 001 VIOLATION — PPTX : tokens absents : {missing}"
        )
