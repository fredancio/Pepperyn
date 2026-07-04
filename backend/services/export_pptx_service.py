"""
export_pptx_service.py — Board Deck Pepperyn
Master de référence officiel : "Ancienne version du Powerpoint.pptx"

Structure IMMUABLE — 16 slides exactes :
  S1   EXECUTIVE DECISION (couverture)
  S2   EXECUTIVE SUMMARY (3 décisions top)
  S3   DIAGNOSTIC (Pourquoi en sommes-nous arrivés là ?)
  S4   GOUVERNANCE / CEO Dashboard (6 KPIs)
  S5   IMPACT FINANCIER (destructeurs de valeur)
  S6   DÉCISION (le coût de l'inaction)
  S7   DÉCISIONS PRIORITAIRES (table)
  S8   EXÉCUTION (roadmap 30/60/90)
  S9   SIMULATION (action vs inaction)
  S10  PROJECTION (trajectoire financière 12 mois)
  S11  RISQUES (risques majeurs)
  S12  PRIORITÉS (matrice impact/effort)
  S13  SUIVI / Carnet d'exécution (table)
  S14  PILOTAGE / Tableau de bord (table)
  S15  LUNDI MATIN (3 décisions clés)
  S16  ANNEXE (qualité données + méthodologie + engagements)

RÈGLE ABSOLUE : seules les données changent. La structure est figée.
Si une donnée numérique est absente → "—". Si un texte est absent → "—".
"""
from __future__ import annotations

import io
import math
import re
from datetime import datetime, timedelta
from typing import Any, Optional, List

# ─── DATE EN FRANÇAIS ─────────────────────────────────────────────────────────
_MOIS_FR = {
    1: "janvier", 2: "février",  3: "mars",     4: "avril",
    5: "mai",     6: "juin",     7: "juillet",  8: "août",
    9: "septembre", 10: "octobre", 11: "novembre", 12: "décembre",
}

def _fr_date(dt: datetime) -> str:
    """Formate une date en français : '02 juillet 2026'."""
    return f"{dt.day:02d} {_MOIS_FR[dt.month]} {dt.year}"

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE

from services.executive_decision_model import build_executive_decision_model
from models.executive_case import ExecutiveCaseJSON

# ─── PALETTE ─────────────────────────────────────────────────────────────────
def _rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

NAVY  = _rgb("0A2540")
BLUE  = _rgb("1B73E8")
RED   = _rgb("C0392B")
AMBER = _rgb("B8763A")
GREEN = _rgb("2C7A4B")
GRAY  = _rgb("8A9BB0")
LGRAY = _rgb("D5DCE5")
WHITE = _rgb("FFFFFF")
DARK  = _rgb("1A1A2E")
LBGRAY = _rgb("F5F7FA")

# ─── DIMENSIONS SLIDE 16:9 (33.87 cm × 19.05 cm) ────────────────────────────
SW = Inches(13.33)
SH = Inches(7.5)
ML = Inches(0.5)    # left margin
MR = Inches(0.5)    # right margin
MT = Inches(1.1)    # top (below header band)
MB = Inches(0.45)   # bottom (above footer)
CW = SW - ML - MR   # content width  ~12.33"
CH = SH - MT - MB   # content height ~5.95"


# ─── HELPERS NUMÉRIQUES ───────────────────────────────────────────────────────

def _fmt_eur(v: Optional[float]) -> str:
    if v is None:
        return "—"
    abs_v = abs(v)
    prefix = "-" if v < 0 else ""
    s = f"{abs_v:,.0f}".replace(",", " ")
    return f"{prefix}{s} €"


def _fmt_auto(v: Optional[float], sign: bool = False) -> str:
    if v is None:
        return "—"
    abs_v = abs(v)
    prefix = "-" if v < 0 else ("+" if sign and v > 0 else "")
    if abs_v >= 950_000:
        m = abs_v / 1_000_000
        if m >= 10:
            return f"{prefix}{m:.0f} M€"
        return f"{prefix}{m:.1f} M€".replace(".", ",")
    return _fmt_eur(v)


def _fmt_chart(v: Optional[float]) -> str:
    """Compact label for chart callout dots — e.g. +620K€, −1,7M€.
    Designed for 0.56"+ wide boxes where the full unit can appear.
    """
    if v is None:
        return "—"
    if v == 0:
        return "0"
    pfx = "−" if v < 0 else "+"
    abs_v = abs(v)
    if abs_v >= 1_000_000:
        m = abs_v / 1_000_000
        s = f"{m:.1f}".replace(".", ",")
        return f"{pfx}{s}M€"
    elif abs_v >= 1_000:
        return f"{pfx}{int(round(abs_v / 1_000))}K€"
    else:
        return f"{pfx}{int(abs_v)}€"


def _fmt_yax(v: Optional[float]) -> str:
    """Y-axis label — compact without € sign to avoid word-wrap in narrow boxes.
    Examples: 898 280 → "+898K" · -1 700 000 → "−1,7M" · 0 → "0"
    """
    if v is None:
        return "—"
    if v == 0:
        return "0"
    pfx = "−" if v < 0 else "+"
    abs_v = abs(v)
    if abs_v >= 1_000_000:
        m = abs_v / 1_000_000
        s = f"{m:.1f}".replace(".", ",")
        return f"{pfx}{s}M"
    elif abs_v >= 1_000:
        return f"{pfx}{int(round(abs_v / 1_000))}K"
    else:
        return f"{pfx}{int(abs_v)}"


def _safe(v, fallback: str = "—") -> str:
    if v is None or v == "" or (isinstance(v, (int, float)) and v == 0):
        return fallback
    return str(v)


def _sm(s) -> str:
    """Strip markdown bold/italic markers (`**`, `*`) from LLM text before PPTX rendering."""
    return re.sub(r'\*+', '', s or '').strip()


# ─── RULE 002 — Helpers anti-troncature ──────────────────────────────────────

def _auto_row_h(text: str, col_w_emu: int, font_pt: int, base_pt: float = 28.0) -> int:
    """
    RULE 002 — Calcule la hauteur minimale en EMU pour contenir ce texte dans
    une colonne de largeur col_w_emu, avec une police de font_pt points.
    Aucune information ne doit être visuellement tronquée.
    """
    col_w_in = col_w_emu / 914400  # EMU → inches
    # Calibri : largeur moyenne d'un caractère ≈ font_pt × 0.52 pt
    chars_per_line = max(1, int((col_w_in * 72) / (font_pt * 0.52)))
    n_lines = max(1, math.ceil(len(str(text)) / chars_per_line))
    # Interligne ≈ 1.35 × font_pt ; padding cellule ≈ 8 pt
    height_pt = n_lines * (font_pt * 1.35) + 8
    return int(Pt(max(base_pt, height_pt)))


def _fit_table_rows(tbl, text_col_idx: int, font_pt: int,
                    skip_header: bool = True, min_h_pt: float = 28.0) -> None:
    """
    RULE 002 — Ajuste les hauteurs de lignes d'un tableau python-pptx pour que
    le texte de la colonne principale (text_col_idx) ne soit jamais tronqué.
    """
    col_w_emu = tbl.columns[text_col_idx].width
    for ri in range(len(tbl.rows)):
        if skip_header and ri == 0:
            continue
        cell = tbl.cell(ri, text_col_idx)
        text = cell.text_frame.text
        needed = _auto_row_h(text, col_w_emu, font_pt, min_h_pt)
        current = tbl.rows[ri].height or 0
        if needed > current:
            tbl.rows[ri].height = needed


# ─── HELPERS PREMIUM : ligne / contexte / barre ──────────────────────────────

def _draw_line_segment(slide, x1: float, y1: float, x2: float, y2: float,
                       color: RGBColor, width_emu: int = None) -> None:
    """Trace un segment de droite via un rectangle fin tourné (Emu coords)."""
    if width_emu is None:
        width_emu = int(Pt(2.5))
    dx, dy = x2 - x1, y2 - y1
    length = math.sqrt(dx * dx + dy * dy)
    if length < 1:
        return
    cx, cy = (x1 + x2) / 2, (y1 + y2) / 2
    shape = slide.shapes.add_shape(
        1,
        int(cx - length / 2), int(cy - width_emu / 2),
        int(length), width_emu,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.rotation = math.degrees(math.atan2(dy, dx))


def _is_growth_context(edm, result: dict) -> bool:
    """True si l'entreprise est stable/en croissance (pas en crise aiguë)."""
    coi = edm.cost_of_inaction
    if coi and coi.per_year and abs(coi.per_year) < 250_000:
        return True
    for card in (result.get("ceo_dashboard") or []):
        if isinstance(card, dict):
            lbl = (card.get("label") or "").upper()
            val = str(card.get("value") or "")
            if "EBITDA" in lbl and val and not val.strip().startswith("-"):
                return True
    return False


def _progress_bar_shape(slide, lx: int, ly: int, w: int, h: int,
                         pct: float, bar_color: RGBColor,
                         bg_color: RGBColor = None) -> None:
    """Barre de progression horizontale (fond + portion remplie)."""
    _rect(slide, lx, ly, w, h, fill_color=bg_color or _rgb("E8ECF2"))
    if pct > 0:
        _rect(slide, lx, ly, max(int(Pt(4)), int(w * min(pct, 1.0))), h,
              fill_color=bar_color)


# ─── UTILITAIRES PYTHON-PPTX ─────────────────────────────────────────────────

def _set_para(para, text: str, size: int, bold: bool = False,
              color: RGBColor = None, align=PP_ALIGN.LEFT,
              italic: bool = False):
    para.text = ""
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = "Calibri"
    if color:
        run.font.color.rgb = color
    para.alignment = align


def _text(slide, txt: str, l, t, w, h,
          size: int = 22, bold: bool = False,
          color: RGBColor = None, align=PP_ALIGN.LEFT,
          italic: bool = False):
    if color is None:
        color = DARK
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    _set_para(para, str(txt), size, bold=bold, color=color, align=align, italic=italic)
    return tb


def _rect(slide, l, t, w, h, fill_color: RGBColor = None, line_color: RGBColor = None):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    return shape


def _header_band(slide, label: str, company: str = ""):
    _rect(slide, 0, 0, SW, Inches(0.75), fill_color=NAVY)
    _text(slide, label.upper(), ML, Inches(0.12), Inches(7), Inches(0.5),
          size=12, bold=True, color=AMBER, align=PP_ALIGN.LEFT)
    _text(slide, "PEPPERYN", SW - Inches(2), Inches(0.12), Inches(1.5), Inches(0.5),
          size=11, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
    _rect(slide, 0, Inches(0.75), SW, Pt(1), fill_color=BLUE)


def _footer_band(slide, page_num: int, company: str = ""):
    fy = SH - Inches(0.4)
    _rect(slide, 0, fy, SW, Inches(0.4), fill_color=_rgb("F0F3F7"))
    _text(slide, "Document confidentiel — usage interne réservé à la direction",
          ML, fy + Pt(4), Inches(9), Inches(0.3),
          size=9, color=GRAY, align=PP_ALIGN.LEFT)
    _text(slide, str(page_num),
          SW - Inches(1), fy + Pt(4), Inches(0.5), Inches(0.3),
          size=9, bold=True, color=GRAY, align=PP_ALIGN.RIGHT)


def _slide_title(slide, title: str, subtitle: str = ""):
    _text(slide, title, ML, MT - Inches(0.25), CW, Inches(0.55),
          size=36, bold=True, color=NAVY, align=PP_ALIGN.LEFT)
    if subtitle:
        _text(slide, subtitle, ML, MT + Inches(0.35), CW, Inches(0.3),
              size=14, color=GRAY, align=PP_ALIGN.LEFT)


# ─── SLIDE 1 : COUVERTURE ─────────────────────────────────────────────────────

def _slide_cover(prs, edm, result: dict, company: str, date_str: str, page: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, SW, SH, fill_color=NAVY)

    # ── Parse date si format dd/mm/yyyy ───────────────────────────────────────
    if date_str and "/" in date_str:
        try:
            parts = date_str.split("/")
            if len(parts) == 3:
                d, m, y = int(parts[0]), int(parts[1]), int(parts[2])
                date_str = f"{d:02d} {_MOIS_FR[m]} {y}"
        except (ValueError, KeyError):
            pass

    # ── Top row ───────────────────────────────────────────────────────────────
    _text(slide, "EXECUTIVE DECISION", Inches(0.5), Inches(0.28), Inches(8), Inches(0.42),
          size=11, bold=True, color=AMBER, align=PP_ALIGN.LEFT)
    _text(slide, "PEPPERYN", SW - Inches(2.6), Inches(0.28), Inches(2.1), Inches(0.42),
          size=11, bold=True, color=_rgb("4A6A8A"), align=PP_ALIGN.RIGHT)
    _rect(slide, Inches(0.5), Inches(0.82), SW - Inches(1.0), Pt(0.5),
          fill_color=_rgb("1B4F8A"))

    # ── Company name ──────────────────────────────────────────────────────────
    _text(slide, company or "—", Inches(0.5), Inches(1.0), Inches(12.3), Inches(1.3),
          size=52, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Date + stats ligne
    n_dec = len(edm.executive_decisions)
    total_impact = sum(d.annual_impact or 0 for d in edm.executive_decisions)
    stats_str = f"{n_dec} décision{'s' if n_dec > 1 else ''} identifiée{'s' if n_dec > 1 else ''}"
    if total_impact:
        stats_str += f"  ·  Impact cumulé : {_fmt_auto(total_impact)} / an"
    _text(slide, date_str, Inches(0.5), Inches(2.38), Inches(8), Inches(0.38),
          size=15, color=_rgb("5A7A9A"), align=PP_ALIGN.LEFT)
    _text(slide, stats_str, Inches(0.5), Inches(2.78), Inches(10), Inches(0.35),
          size=11, color=_rgb("3A6A8A"), align=PP_ALIGN.LEFT)

    # ── Decision — left column ────────────────────────────────────────────────
    dec = edm.executive_decisions[0] if edm.executive_decisions else None
    dec_txt = _sm(dec.decision) if dec else "Décision prioritaire non disponible"
    impact_str = (_fmt_auto(dec.annual_impact) + " / an" if dec and dec.annual_impact
                  else "Non chiffrable")

    _text(slide, "DÉCISION PRIORITAIRE #1", Inches(0.5), Inches(2.95),
          Inches(8.5), Inches(0.38),
          size=9, bold=True, color=AMBER, align=PP_ALIGN.LEFT)
    _text(slide, dec_txt, Inches(0.5), Inches(3.4), Inches(8.5), Inches(1.55),
          size=22, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # ── Impact box — right column ─────────────────────────────────────────────
    box_w = Inches(3.7)
    box_h = Inches(1.45)
    box_x = SW - box_w - Inches(0.45)
    box_y = Inches(2.85)
    impact_bg = slide.shapes.add_shape(1, box_x, box_y, box_w, box_h)
    impact_bg.fill.solid()
    impact_bg.fill.fore_color.rgb = _rgb("061828")
    impact_bg.line.color.rgb = AMBER
    impact_bg.line.width = Pt(1.5)
    _text(slide, impact_str, box_x + Inches(0.15), box_y + Inches(0.1),
          box_w - Inches(0.3), Inches(0.9),
          size=34, bold=True, color=AMBER, align=PP_ALIGN.CENTER)
    _text(slide, "IMPACT ANNUEL ESTIMÉ", box_x + Inches(0.15), box_y + Inches(0.98),
          box_w - Inches(0.3), Inches(0.32),
          size=8, color=_rgb("6A8AA0"), align=PP_ALIGN.CENTER)

    # ── Bottom bar ────────────────────────────────────────────────────────────
    _text(slide, "RAPPORT CONFIDENTIEL — USAGE INTERNE",
          Inches(0.5), SH - Inches(0.48), Inches(10), Inches(0.35),
          size=9, color=_rgb("3A4A5A"), align=PP_ALIGN.LEFT)
    _text(slide, str(page), SW - Inches(1), SH - Inches(0.48), Inches(0.5), Inches(0.35),
          size=9, color=GRAY, align=PP_ALIGN.RIGHT)


# ─── SLIDE 2 : EXECUTIVE SUMMARY ─────────────────────────────────────────────

def _slide_exec_summary(prs, edm, result: dict, company: str, date_str: str, page: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header_band(slide, "EXECUTIVE SUMMARY", company)
    _footer_band(slide, page, company)
    _slide_title(slide, "Synthèse des décisions prioritaires")

    decisions = edm.executive_decisions[:3]
    if not decisions:
        _text(slide, "—", ML, MT + Inches(0.5), CW, Inches(1), size=22, color=GRAY)
        return

    growth = _is_growth_context(edm, result)
    # Cartes compactes (was ~1.8" each, now 1.42") → libère de l'espace pour la ligne total
    row_h = Inches(1.42)
    row_top = MT + Inches(0.55)

    for i, dec in enumerate(decisions):
        y = row_top + i * (row_h + Pt(4))
        card_color = [BLUE, _rgb("2C7A4B"), _rgb("B8763A")][i]
        _rect(slide, ML, y, CW, row_h, fill_color=_rgb("F0F4FA"), line_color=LGRAY)
        # Accent gauche coloré
        _rect(slide, ML, y, Inches(0.07), row_h, fill_color=card_color)
        _text(slide, f"#{i + 1}", ML + Inches(0.18), y + Inches(0.08), Inches(0.4), row_h,
              size=16, bold=True, color=card_color)
        _text(slide, _sm(dec.decision), ML + Inches(0.65), y + Inches(0.08),
              Inches(5.5), row_h - Inches(0.15), size=14, color=DARK)
        # Propriétaire
        owner_str = dec.owner or "Direction"
        _text(slide, owner_str, ML + Inches(0.65), y + Inches(0.9),
              Inches(3), Inches(0.3), size=10, color=GRAY, italic=True)
        impact_str = _fmt_auto(dec.annual_impact, sign=True) if dec.annual_impact else "—"
        impact_color = RED if not growth else GREEN
        _text(slide, impact_str, ML + Inches(6.4), y + Inches(0.08), Inches(1.9), Inches(0.5),
              size=17, bold=True, color=impact_color, align=PP_ALIGN.RIGHT)
        _text(slide, "IMPACT / AN", ML + Inches(6.4), y + Inches(0.6), Inches(1.9), Inches(0.28),
              size=9, color=GRAY, align=PP_ALIGN.RIGHT)
        roi = f"{dec.roi_score:.1f}/10" if dec.roi_score else "—"
        _text(slide, roi, ML + Inches(8.5), y + Inches(0.08), Inches(1.2), Inches(0.5),
              size=17, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        _text(slide, "SCORE ROI", ML + Inches(8.5), y + Inches(0.6), Inches(1.2), Inches(0.28),
              size=9, color=GRAY, align=PP_ALIGN.CENTER)
        horizon = dec.timeline or "—"
        _text(slide, horizon, ML + Inches(9.85), y + Inches(0.08), Inches(1.2), Inches(0.5),
              size=14, color=AMBER, align=PP_ALIGN.CENTER)
        _text(slide, "HORIZON", ML + Inches(9.85), y + Inches(0.6), Inches(1.2), Inches(0.28),
              size=9, color=GRAY, align=PP_ALIGN.CENTER)
        status = dec.status or "À lancer"
        _text(slide, status, ML + Inches(11.2), y + Inches(0.08), Inches(1.1), Inches(0.5),
              size=12, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
        _text(slide, "STATUT", ML + Inches(11.2), y + Inches(0.6), Inches(1.1), Inches(0.28),
              size=9, color=GRAY, align=PP_ALIGN.CENTER)

    # ── Ligne impact total ─────────────────────────────────────────────────────
    total_impact = sum(d.annual_impact or 0 for d in decisions)
    if total_impact:
        ty = int(row_top + 3 * (row_h + Pt(4)) + Inches(0.14))
        _rect(slide, int(ML), ty, int(CW), int(Inches(0.52)), fill_color=NAVY)
        _text(slide, "IMPACT TOTAL CUMULÉ — SI LES 3 DÉCISIONS SONT ENGAGÉES",
              int(ML + Inches(0.2)), ty + int(Inches(0.09)),
              int(Inches(7)), int(Inches(0.34)),
              size=10, bold=True, color=GRAY, align=PP_ALIGN.LEFT)
        _text(slide, _fmt_auto(total_impact, sign=True),
              int(ML + Inches(7.5)), ty + int(Inches(0.06)),
              int(Inches(3.5)), int(Inches(0.4)),
              size=16, bold=True, color=AMBER, align=PP_ALIGN.RIGHT)
        _text(slide, "/ AN",
              int(ML + Inches(11.1)), ty + int(Inches(0.12)),
              int(Inches(1)), int(Inches(0.28)),
              size=10, color=GRAY, align=PP_ALIGN.LEFT)


# ─── SLIDE 3 : DIAGNOSTIC ─────────────────────────────────────────────────────

def _slide_diagnostic(prs, edm, result: dict, company: str, date_str: str, page: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header_band(slide, "DIAGNOSTIC", company)
    _footer_band(slide, page, company)
    growth = _is_growth_context(edm, result)
    title = ("Où capturer votre prochaine vague de valeur"
             if growth else "Pourquoi en sommes-nous arrivés là ?")
    _slide_title(slide, title)

    diag = (result.get("diagnostic_immediat") or result.get("resume_executif")
            or result.get("synthese") or "—")
    diag = diag.strip()
    if len(diag) > 350:
        diag = diag[:347] + "…"
    _text(slide, diag, ML, MT + Inches(0.55), CW, Inches(1.1), size=14, color=DARK)

    tension = result.get("phrase_tension")
    if tension:
        tp_top = MT + Inches(1.73)
        tp_bg = slide.shapes.add_shape(1, ML, tp_top, CW, Inches(0.52))
        tp_bg.fill.solid()
        tp_bg.fill.fore_color.rgb = _rgb("0A2540")
        tp_bg.line.fill.background()
        _text(slide, tension.strip(), ML + Inches(0.15), tp_top + Pt(5),
              CW - Inches(0.3), Inches(0.46), size=13, bold=True, color=WHITE, italic=True)

    destroyers = edm.value_destroyers[:3]
    lbl = "TOP 3 LEVIERS D'OPTIMISATION" if not growth else "TOP 3 LEVIERS DE CRÉATION DE VALEUR"
    if destroyers:
        _text(slide, lbl, ML, MT + Inches(2.37), CW, Inches(0.38),
              size=12, bold=True, color=BLUE)
        card_w = CW / 3 - Inches(0.12)
        card_h = Inches(2.85)
        for i, d in enumerate(destroyers):
            lx = ML + i * (card_w + Inches(0.18))
            ly = MT + Inches(2.78)
            bg   = _rgb("EFF8F1") if growth else _rgb("FBF0EE")
            bord = _rgb("8EC9A2") if growth else _rgb("E8B0A8")
            val_c = GREEN if growth else RED
            _rect(slide, lx, ly, card_w, card_h, fill_color=bg, line_color=bord)
            # M1 fix: afficher abs() — les leviers sont toujours des opportunités positives
            ann_str = _fmt_auto(abs(d.annual_impact)) if d.annual_impact else "Non chiffré"
            _text(slide, ann_str, lx + Inches(0.1), ly + Inches(0.15),
                  card_w - Inches(0.2), Inches(0.7), size=24, bold=True, color=val_c)
            _text(slide, "/ AN", lx + Inches(0.1), ly + Inches(0.82),
                  card_w - Inches(0.2), Inches(0.25), size=9, color=GRAY)
            _text(slide, _sm(d.name), lx + Inches(0.1), ly + Inches(1.1),
                  card_w - Inches(0.2), Inches(1.3), size=13, bold=True, color=DARK)
            # Mini waterfall bar visuel
            if d.annual_impact:
                max_abs = max(abs(x.annual_impact or 0) for x in destroyers) or 1
                bar_pct = abs(d.annual_impact) / max_abs
                bar_w = int((card_w - Inches(0.2)) * bar_pct)
                _rect(slide, int(lx + Inches(0.1)), int(ly + card_h - Inches(0.35)),
                      bar_w, int(Inches(0.12)), fill_color=val_c)


# ─── SLIDE 4 : GOUVERNANCE / CEO DASHBOARD ────────────────────────────────────

def _slide_dashboard(prs, edm, result: dict, company: str, date_str: str, page: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header_band(slide, "GOUVERNANCE", company)
    _footer_band(slide, page, company)
    _slide_title(slide, "CEO Dashboard", "Indicateurs stratégiques clés")

    dashboard = result.get("ceo_dashboard") or []
    items = []
    for card in dashboard[:6]:
        if isinstance(card, dict):
            items.append(card)
        else:
            items.append({"label": getattr(card, "label", ""),
                          "value": getattr(card, "value", ""),
                          "status": getattr(card, "status", None)})
    while len(items) < 6:
        items.append({"label": "—", "value": "—", "status": "missing"})

    card_w = CW / 3 - Inches(0.12)
    card_h = Inches(1.95)
    card_top_1 = MT + Inches(0.65)
    card_top_2 = card_top_1 + card_h + Inches(0.18)

    for idx, item in enumerate(items[:6]):
        row = idx // 3
        col = idx % 3
        lx = ML + col * (card_w + Inches(0.18))
        ly = card_top_1 if row == 0 else card_top_2
        is_missing = (item.get("status") == "missing" or
                      "données insuf" in str(item.get("value", "")).lower())
        border_c = LGRAY if is_missing else BLUE
        val_c = GRAY if is_missing else RED
        val = str(item.get("value", "—"))
        lbl = str(item.get("label", ""))
        _rect(slide, lx, ly, card_w, card_h, fill_color=_rgb("F5F8FF"), line_color=border_c)
        val_size = 24 if len(val) <= 22 else (18 if len(val) <= 45 else 13)
        _text(slide, val, lx + Inches(0.12), ly + Inches(0.25),
              card_w - Inches(0.24), Inches(0.85),
              size=val_size, bold=True, color=val_c, align=PP_ALIGN.CENTER)
        _text(slide, lbl.upper(), lx + Inches(0.1), ly + Inches(1.2),
              card_w - Inches(0.2), Inches(0.5),
              size=11, bold=True, color=GRAY, align=PP_ALIGN.CENTER)


# ─── SLIDE 5 : IMPACT FINANCIER ───────────────────────────────────────────────

def _slide_impact_financier(prs, edm, result: dict, company: str, date_str: str, page: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header_band(slide, "IMPACT FINANCIER", company)
    _footer_band(slide, page, company)
    _slide_title(slide, "Où la valeur attend d'être libérée", "Analyse des leviers d'optimisation")

    destroyers = edm.value_destroyers[:8]
    rows = []
    for d in destroyers:
        rows.append([
            _sm(d.name),
            _fmt_auto(d.annual_impact) if d.annual_impact else "Non chiffrable",
            _fmt_eur(d.monthly_impact) if d.monthly_impact else "—",
            d.trend or "—",
        ])
    if not rows:
        rows = [["—", "—", "—", "—"]]

    tbl_top = MT + Inches(0.75)      # subtitle bottom ≈ MT + 0.65" → safe margin
    hdr_h   = int(Inches(0.50))
    row_h   = int(Inches(0.58))
    tbl_h   = hdr_h + len(rows) * row_h
    tbl = slide.shapes.add_table(len(rows) + 1, 4, ML, tbl_top, CW, tbl_h).table
    tbl.rows[0].height = hdr_h
    for _ri in range(1, len(rows) + 1):
        tbl.rows[_ri].height = row_h

    col_pct = [0.48, 0.22, 0.18, 0.12]
    for i, pct in enumerate(col_pct):
        tbl.columns[i].width = int(CW * pct)

    for ci, hdr in enumerate(["Destructeur de valeur", "Impact annuel", "Impact mensuel", "Tendance"]):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        _set_para(cell.text_frame.paragraphs[0], hdr, 11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        cell.margin_top = Pt(4)
        cell.margin_bottom = Pt(4)

    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            if ri % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb("F5F8FF")
            else:
                cell.fill.background()
            fc = RED if ci == 1 else DARK
            _set_para(cell.text_frame.paragraphs[0], val, 12,
                      bold=(ci == 1), color=fc,
                      align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER)
    _fit_table_rows(tbl, text_col_idx=0, font_pt=12)  # RULE 002

    # ── Ligne total ───────────────────────────────────────────────────────────
    # C4 fix: utiliser abs() pour chaque levier → total = COI réel (161K)
    # et non la somme algébrique (+119K - 42K = 77K) qui sous-estime le potentiel.
    total_annual = sum(abs(d.annual_impact or 0) for d in destroyers)
    total_monthly = sum(abs(d.monthly_impact or 0) for d in destroyers)
    if total_annual:
        ty = int(tbl_top + tbl_h + Inches(0.08))
        _rect(slide, int(ML), ty, int(CW), int(Inches(0.48)), fill_color=NAVY)
        # mn1 fix: label précis → "Coût de l'inaction" (161K = ce qu'on perd),
        # distinct de "Impact total décisions" (139K = ce qu'on capture via D1+D2+D3).
        _text(slide, "TOTAL — Coût de l'inaction · par an",
              int(ML + Inches(0.15)), ty + int(Inches(0.08)),
              int(CW * 0.48), int(Inches(0.32)),
              size=11, bold=True, color=WHITE)
        _text(slide, _fmt_auto(total_annual),
              int(ML + CW * 0.48), ty + int(Inches(0.05)),
              int(CW * 0.22), int(Inches(0.38)),
              size=14, bold=True, color=AMBER, align=PP_ALIGN.CENTER)
        _text(slide, _fmt_eur(total_monthly) if total_monthly else "—",
              int(ML + CW * 0.70), ty + int(Inches(0.05)),
              int(CW * 0.18), int(Inches(0.38)),
              size=11, bold=True, color=_rgb("C8A050"), align=PP_ALIGN.CENTER)


# ─── SLIDE 6 : DÉCISION — COÛT DE L'INACTION ─────────────────────────────────

def _slide_cout_inaction(prs, edm, result: dict, company: str, date_str: str, page: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    growth = _is_growth_context(edm, result)
    _header_band(slide, "DÉCISION", company)
    _footer_band(slide, page, company)
    if growth:
        _slide_title(slide, "La valeur qui n'attend pas", "Chaque trimestre sans action = valeur non capturée")
    else:
        _slide_title(slide, "Le coût de l'inaction", "Si rien ne change maintenant")

    coi = edm.cost_of_inaction
    hero_str = (_fmt_auto(abs(coi.per_year)) if coi and coi.per_year else "—")
    hero_color = GREEN if growth else RED
    label_txt = "VALEUR NON CAPTURÉE — PAR AN" if growth else "COÛT DE L'INACTION — PAR AN"
    _text(slide, hero_str, ML, MT + Inches(0.45), Inches(8), Inches(1.1),
          size=56, bold=True, color=hero_color, align=PP_ALIGN.LEFT)
    _text(slide, label_txt, ML, MT + Inches(1.62),
          Inches(8), Inches(0.35), size=13, bold=True, color=hero_color)
    _rect(slide, ML, MT + Inches(2.05), Inches(8), Pt(2), fill_color=hero_color)

    sub_vals = [
        (_fmt_eur(abs(coi.per_month)) if coi and coi.per_month else "—", "PAR MOIS"),
        (_fmt_eur(abs(coi.per_week))  if coi and coi.per_week  else "—", "PAR SEMAINE"),
        (_fmt_eur(abs(coi.per_day))   if coi and coi.per_day   else "—", "PAR JOUR"),
    ]
    sub_w = Inches(3.5)
    for i, (val, lbl) in enumerate(sub_vals):
        lx = ML + i * (sub_w + Inches(0.4))
        _text(slide, val, lx, MT + Inches(2.28), sub_w, Inches(0.65),
              size=24, bold=True, color=DARK)
        _text(slide, lbl, lx, MT + Inches(2.9), sub_w, Inches(0.3),
              size=10, color=GRAY)

    risque = result.get("risque_inaction") or ""
    # A fix: remplacer les montants approximatifs LLM (~136 K€ etc.) par la valeur EDM réelle (COI/an).
    if risque and coi and coi.per_year:
        import re as _re_r
        _coi_ke = f"{int(round(abs(coi.per_year) / 1_000))} K€"
        risque = _re_r.sub(r'~\s*\d[\d\s]*\s*K€', f'~{_coi_ke}', risque)
    if risque:
        _rect(slide, ML, MT + Inches(3.35), CW, Inches(0.72),
              fill_color=_rgb("F5F8FF"), line_color=LGRAY)
        _text(slide, risque, ML + Inches(0.15), MT + Inches(3.44),
              CW - Inches(0.3), Inches(0.62), size=13, color=DARK, italic=True)

    # ── Timeline visuelle : semaines restantes ────────────────────────────────
    # Estimé à partir du COI mensuel vs trésorerie (si disponible)
    weeks = None
    for card in (result.get("ceo_dashboard") or []):
        if isinstance(card, dict) and "TRÉSORERIE" in (card.get("label") or "").upper():
            try:
                cash_str = str(card.get("value") or "").replace(" ", "").replace("€", "")
                cash_str = cash_str.replace("K", "000").replace("M", "000000")
                cash_val = float(cash_str.replace(",", "."))
                if coi and coi.per_month and coi.per_month != 0:
                    weeks = max(1, int(cash_val / abs(coi.per_month) * 4))
            except (ValueError, TypeError):
                pass

    if weeks is not None and not growth:
        ty = MT + Inches(4.22)
        _text(slide, f"À trésorerie constante : point critique estimé dans {weeks} semaine{'s' if weeks > 1 else ''}",
              ML, ty, CW, Inches(0.32), size=11, bold=True, color=_rgb("CC3333"))
        # Barre de progression (countdown visuel)
        BAR_W = int(CW)
        BAR_H = int(Inches(0.22))
        bar_top = int(ty + Inches(0.38))
        _rect(slide, int(ML), bar_top, BAR_W, BAR_H, fill_color=_rgb("E8ECF2"))
        # Portion "déjà consommée" en rouge (80% par défaut si < 12 semaines)
        urgency_pct = max(0.0, min(1.0, 1.0 - weeks / 52.0))
        if urgency_pct > 0:
            _rect(slide, int(ML), bar_top, int(BAR_W * urgency_pct), BAR_H,
                  fill_color=_rgb("CC3333"))
        _text(slide, "MAINTENANT", int(ML), bar_top + int(Inches(0.28)),
              int(Inches(2)), int(Inches(0.22)), size=8, color=GRAY)
        _text(slide, f"J+{weeks * 7}j", int(ML + BAR_W - Inches(1.5)), bar_top + int(Inches(0.28)),
              int(Inches(1.5)), int(Inches(0.22)), size=8, bold=True, color=_rgb("CC3333"),
              align=PP_ALIGN.RIGHT)


# ─── SLIDE 7 : DÉCISIONS PRIORITAIRES ────────────────────────────────────────

def _slide_decisions_prioritaires(prs, edm, result: dict, company: str, date_str: str, page: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header_band(slide, "DÉCISIONS PRIORITAIRES", company)
    _footer_band(slide, page, company)
    _slide_title(slide, "Plan d'action prioritaire")

    # ── EDX-001 : récupération du raisonnement ────────────────────────────────
    reasoning_list = result.get("decision_reasoning", []) or []
    has_reasoning = any(
        r.get("why_this_decision") or r.get("problem_source")
        for r in reasoning_list
    )
    reasoning_by_idx = {r["decision_index"]: r for r in reasoning_list}

    # CV-C03 — panneau agrandi pour raisonnement lisible (min 10pt)
    REASONING_PANEL_H = Inches(2.5) if has_reasoning else Inches(0)

    decisions = edm.executive_decisions[:8]
    rows = []
    for d in decisions:
        rows.append([
            _sm(d.decision),
            _fmt_auto(d.annual_impact, sign=True) if d.annual_impact else "Non chiffrable",
            f"{d.roi_score:.1f}/10",
            d.priority or "—",
            d.owner or "Direction",
            d.timeline or "—",
            d.status or "À lancer",
        ])
    if not rows:
        rows = [["—"] + ["—"] * 6]

    col_pct     = [0.26, 0.14, 0.07, 0.10, 0.13, 0.14, 0.16]
    tbl_top     = MT + Inches(0.55)
    available_h = CH - Inches(0.6) - REASONING_PANEL_H
    hdr_h       = int(Inches(0.45))
    max_row_h   = int(min(Inches(0.55), (available_h - hdr_h) / max(len(rows), 1)))
    row_h       = max(int(Inches(0.35)), max_row_h)
    tbl_h       = hdr_h + len(rows) * row_h
    tbl = slide.shapes.add_table(len(rows) + 1, 7, ML, tbl_top, CW, tbl_h).table
    tbl.rows[0].height = hdr_h
    for _ri in range(1, len(rows) + 1):
        tbl.rows[_ri].height = row_h
    for i, pct in enumerate(col_pct):
        tbl.columns[i].width = int(CW * pct)

    for ci, hdr in enumerate(["Décision", "Impact annuel", "ROI", "Priorité", "Responsable", "Horizon", "Statut"]):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        _set_para(cell.text_frame.paragraphs[0], hdr, 10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            if ri % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb("F5F8FF")
            else:
                cell.fill.background()
            fc = RED if ci == 1 else (AMBER if ci == 3 else DARK)
            _set_para(cell.text_frame.paragraphs[0], val, 10, color=fc,
                      align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER)
    _fit_table_rows(tbl, text_col_idx=0, font_pt=10, min_h_pt=32.0)  # RULE 002

    # ── EDX-001 : panneau chaîne décisionnelle ────────────────────────────────
    # Affiché sous la table, dans la même slide, sans nouvelle slide.
    if has_reasoning:
        panel_top = tbl_top + tbl_h + Inches(0.12)
        panel_h = REASONING_PANEL_H - Inches(0.12)

        # Fond légèrement teinté
        panel_bg = slide.shapes.add_shape(
            1, ML, panel_top, CW, panel_h  # MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE=5; RECTANGLE=1
        )
        panel_bg.fill.solid()
        panel_bg.fill.fore_color.rgb = _rgb("EEF3FB")
        panel_bg.line.color.rgb = _rgb("C5D5F0")
        panel_bg.line.width = Pt(0.5)

        # Titre du panneau
        title_box = slide.shapes.add_textbox(
            ML + Inches(0.08), panel_top + Inches(0.06),
            CW - Inches(0.16), Inches(0.22)
        )
        tf = title_box.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "RAISONNEMENT — POURQUOI CES DÉCISIONS"
        run.font.size = Pt(9)
        run.font.bold = True
        run.font.color.rgb = _rgb("1B73E8")

        # CV-C03 — contenu lisible : top 3 décisions, 2 lignes par décision, 10pt
        content_top = panel_top + Inches(0.30)
        content_h = panel_h - Inches(0.34)
        content_box = slide.shapes.add_textbox(
            ML + Inches(0.08), content_top,
            CW - Inches(0.16), content_h
        )
        tf2 = content_box.text_frame
        tf2.word_wrap = True

        shown = 0
        for i, dec in enumerate(decisions[:3]):  # Top 3 uniquement — lisibilité
            r = reasoning_by_idx.get(i, {})
            why = r.get("why_this_decision")
            prob = r.get("problem_source")
            conf = r.get("decision_confidence")
            if not why and not prob:
                continue

            # Ligne 1 : numéro + décision (gras)
            para1 = tf2.paragraphs[0] if shown == 0 else tf2.add_paragraph()
            run_title = para1.add_run()
            run_title.text = f"#{i+1}  {_sm(dec.decision)}"
            run_title.font.size = Pt(10)
            run_title.font.bold = True
            run_title.font.color.rgb = _rgb("0A2540")

            # Ligne 2 : raisonnement (normal)
            para2 = tf2.add_paragraph()
            run_why = para2.add_run()
            why_text = why or prob or ""
            # Tronquer à 120 caractères pour tenir sur une ligne
            if len(why_text) > 120:
                why_text = why_text[:117] + "…"
            conf_tag = f"  [Confiance : {conf}%]" if conf is not None else ""
            run_why.text = f"→ {why_text}{conf_tag}"
            run_why.font.size = Pt(10)
            run_why.font.color.rgb = _rgb("444444")

            # Séparateur vide entre décisions
            if shown < 2:
                tf2.add_paragraph()

            shown += 1


# ─── SLIDE 7b : RAISONNEMENT — MÉTHODOLOGIE (fallback) ───────────────────────

def _slide_methodologie(prs, edm, result: dict, company: str, date_str: str, page: int):
    """
    Slide de raisonnement stratégique — POURQUOI ces décisions, pas les alternatives.
    Fallback utilisé quand pas de données EDX-002 (options_considered absentes).
    La slide répond à la question du Codir : "Avez-vous étudié d'autres pistes ?"
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header_band(slide, "RAISONNEMENT STRATÉGIQUE", company)
    _footer_band(slide, page, company)
    _slide_title(slide, "Pourquoi ces décisions, et pas d'autres",
                 "Critères de sélection : impact mesuré · ROI · urgence · faisabilité opérationnelle")

    decisions = edm.executive_decisions[:4]
    reasoning_list = result.get("decision_reasoning", []) or []

    if not decisions:
        _text(slide, "—", ML, MT + Inches(0.75), CW, Inches(1), size=22, color=GRAY)
        return

    # Layout : 2×2 (4 décisions) ou 1×3 (3 décisions) ou 1×2
    n = len(decisions)
    if n >= 4:
        card_w = CW / 2 - Inches(0.12)
        card_h = (CH - Inches(0.85)) / 2 - Inches(0.08)
        positions = [
            (ML,                           MT + Inches(0.75)),
            (ML + card_w + Inches(0.24),   MT + Inches(0.75)),
            (ML,                           MT + Inches(0.75) + card_h + Inches(0.16)),
            (ML + card_w + Inches(0.24),   MT + Inches(0.75) + card_h + Inches(0.16)),
        ]
    elif n == 3:
        card_w = CW / 3 - Inches(0.12)
        card_h = CH - Inches(0.85)
        positions = [(ML + i * (card_w + Inches(0.18)), MT + Inches(0.75)) for i in range(3)]
    else:
        card_w = CW / 2 - Inches(0.12)
        card_h = CH - Inches(0.85)
        positions = [(ML + i * (card_w + Inches(0.24)), MT + Inches(0.75)) for i in range(n)]

    colors = [BLUE, AMBER, GREEN, RED]

    for i, dec in enumerate(decisions[:4]):
        if i >= len(positions):
            break
        lx, ly = int(positions[i][0]), int(positions[i][1])
        c = colors[i]

        # Fond carte
        _rect(slide, lx, ly, int(card_w), int(card_h), fill_color=_rgb("F5F8FF"), line_color=c)
        _rect(slide, lx, ly, int(card_w), int(Inches(0.055)), fill_color=c)

        # Numéro + décision
        _text(slide, f"#{i+1}", lx + int(Inches(0.1)), ly + int(Inches(0.1)),
              int(Inches(0.35)), int(Inches(0.5)), size=22, bold=True, color=c)
        _text(slide, _sm(dec.decision), lx + int(Inches(0.5)), ly + int(Inches(0.1)),
              int(card_w - Inches(0.62)), int(min(Inches(1.0), card_h * 0.38)),
              size=10, bold=True, color=NAVY)

        # Impact + ROI
        impact_str = _fmt_auto(dec.annual_impact) if dec.annual_impact else "—"
        roi_str = f"ROI {dec.roi_score:.1f}/10" if dec.roi_score else ""
        _text(slide, impact_str, lx + int(Inches(0.1)), ly + int(card_h * 0.42),
              int(card_w - Inches(0.2)), int(Inches(0.42)), size=15, bold=True, color=RED)
        if roi_str:
            _text(slide, f"{roi_str}  ·  {dec.owner or 'Direction'}  ·  {dec.timeline or '—'}",
                  lx + int(Inches(0.1)), ly + int(card_h * 0.42) + int(Inches(0.44)),
                  int(card_w - Inches(0.2)), int(Inches(0.28)), size=9, color=GRAY)

        # Raisonnement
        r = next((x for x in reasoning_list if x.get("decision_index") == i), {})
        why = r.get("why_this_decision") or r.get("problem_source") or ""
        if not why and dec.decision:
            why = f"Cette décision adresse directement le principal levier d'optimisation identifié dans les données."
        if why:
            if len(why) > 180:
                why = why[:177] + "…"
            sep_y = ly + int(card_h * 0.42) + int(Inches(0.72))
            sep = slide.shapes.add_shape(1, lx + int(Inches(0.1)), sep_y,
                                         int(card_w - Inches(0.2)), int(Pt(0.5)))
            sep.fill.solid()
            sep.fill.fore_color.rgb = c
            sep.line.fill.background()
            _text(slide, f"→ {why}",
                  lx + int(Inches(0.1)), sep_y + int(Inches(0.1)),
                  int(card_w - Inches(0.2)), int(card_h - (sep_y - ly) - Inches(0.18)),
                  size=9, color=DARK, italic=True)



# ─── SLIDE 7b : RAISONNEMENT DÉCISIONNEL (EDX-002) ───────────────────────────
# La slide McKinsey. "J'ai étudié N options. J'en ai éliminé N-1. Voici pourquoi."
# Si pas de données EDX-002 → slide méthodologie (toujours affichée).

def _slide_raisonnement_comparatif(prs, edm, result: dict, company: str, date_str: str, page: int):
    """
    Slide de raisonnement comparatif EDX-002.
    Montre pour les 2 premières décisions :
    - Les options évaluées et leur critère d'élimination
    - Pourquoi l'option retenue domine
    - Ce qui ferait changer d'avis (conditions de révision)
    Fallback : slide méthodologie 4 phases (toujours une slide créée).
    """
    reasoning_list = result.get("decision_reasoning", []) or []
    decisions_with_options = [
        r for r in reasoning_list
        if r.get("options_considered")
    ]

    # Fallback systématique : slide de méthodologie si pas de données EDX-002
    # (garantit que la slide existe toujours → numérotation footer correcte)
    if not decisions_with_options:
        _slide_methodologie(prs, edm, result, company, date_str, page)
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header_band(slide, "RAISONNEMENT DÉCISIONNEL", company)
    _footer_band(slide, page, company)
    _slide_title(slide, "Comment Pepperyn a raisonné",
                 "Options évaluées • Éliminations • Décision dominante • Conditions de révision")

    decisions = edm.executive_decisions
    reasoning_by_idx = {r["decision_index"]: r for r in reasoning_list}

    # Layout : jusqu'à 2 décisions côte à côte ou en colonne selon la disponibilité
    top_decisions = [r for r in decisions_with_options[:2]]
    n_blocks = len(top_decisions)

    # Largeur et hauteur de chaque bloc
    block_w = CW if n_blocks == 1 else (CW - Inches(0.15)) / 2
    block_h = CH - Inches(0.82)   # shorter to compensate for later block_top
    block_top = MT + Inches(0.72)  # subtitle ends at MT+0.65" → start after with margin

    for b_idx, r in enumerate(top_decisions):
        bx = ML + b_idx * (block_w + Inches(0.15))
        dec_idx = r["decision_index"]
        dec = next((d for i, d in enumerate(decisions) if i == dec_idx), None)
        options   = r.get("options_considered") or []
        dominant  = r.get("dominant_rationale") or ""
        tippings  = r.get("tipping_conditions") or []

        # Fond du bloc
        bg = slide.shapes.add_shape(1, bx, block_top, block_w, block_h)
        bg.fill.solid()
        bg.fill.fore_color.rgb = _rgb("F5F8FF")
        bg.line.color.rgb = _rgb("C5D5F0")
        bg.line.width = Pt(0.5)

        # Titre de la décision (strip markdown bold markers from LLM output)
        dec_name = re.sub(r'\*+', '', dec.decision or '').strip() if dec else f"Décision #{dec_idx + 1}"
        impact_str = (_fmt_auto(dec.annual_impact) if dec and dec.annual_impact else "")
        header_text = f"#{dec_idx + 1} — {dec_name}"
        if impact_str:
            header_text += f"  |  {impact_str}/an"
        _text(slide, header_text, bx + Inches(0.1), block_top + Inches(0.08),
              block_w - Inches(0.2), Inches(0.32),
              size=11, bold=True, color=NAVY)

        # Sous-titre "N options évaluées"
        n_opt = len(options)
        n_elim = n_opt
        subtitle_text = (
            f"Pepperyn a évalué {n_opt + 1} option{'s' if n_opt > 0 else ''} — "
            f"{n_elim} écartée{'s' if n_elim != 1 else ''}"
        )
        _text(slide, subtitle_text, bx + Inches(0.1), block_top + Inches(0.38),
              block_w - Inches(0.2), Inches(0.2),
              size=9, color=_rgb("1B73E8"), italic=True)

        # Séparateur
        sep = slide.shapes.add_shape(1, bx + Inches(0.1),
                                      block_top + Inches(0.56),
                                      block_w - Inches(0.2), Pt(0.5))
        sep.fill.solid()
        sep.fill.fore_color.rgb = _rgb("C5D5F0")
        sep.line.fill.background()

        # Options éliminées
        opts_top = block_top + Inches(0.64)
        opts_h = Inches(0.26)
        for opt in options[:3]:  # Max 3 options pour lisibilité et aération
            opt_name = opt.get("option", "")[:70]
            elim = opt.get("elimination_criterion", "")[:90]
            if len(opt.get("option", "")) > 70:
                opt_name += "…"
            if len(opt.get("elimination_criterion", "")) > 90:
                elim += "…"
            _text(slide, f"✗  {opt_name}",
                  bx + Inches(0.12), opts_top,
                  block_w - Inches(0.24), opts_h,
                  size=9, bold=False, color=_rgb("CC3333"))
            opts_top += opts_h
            if elim:
                _text(slide, f"   → {elim}",
                      bx + Inches(0.12), opts_top,
                      block_w - Inches(0.24), opts_h,
                      size=8, color=_rgb("777777"), italic=True)
                opts_top += opts_h

        # Option dominante (si espace disponible)
        if dominant and opts_top < block_top + block_h - Inches(0.75):
            # Fond distinct pour l'option dominante
            dom_top = opts_top + Inches(0.08)
            dom_h = Inches(0.65)  # was 0.35 — hauteur pour 2 lignes à 9pt
            dom_bg = slide.shapes.add_shape(1, bx + Inches(0.1), dom_top,
                                            block_w - Inches(0.2), dom_h)
            dom_bg.fill.solid()
            dom_bg.fill.fore_color.rgb = _rgb("0A2540")
            dom_bg.line.fill.background()

            dominant_short = dominant[:200] + "…" if len(dominant) > 200 else dominant
            _text(slide, f"✓  {dominant_short}",
                  bx + Inches(0.18), dom_top + Inches(0.06),
                  block_w - Inches(0.36), dom_h - Inches(0.10),
                  size=9, bold=True, color=WHITE)
            opts_top = dom_top + dom_h

        # Conditions de révision (si espace)
        if tippings and opts_top < block_top + block_h - Inches(0.25):
            tip_top = opts_top + Inches(0.1)
            _text(slide, "Conditions de révision :",
                  bx + Inches(0.12), tip_top,
                  block_w - Inches(0.24), Inches(0.18),
                  size=8, bold=True, color=AMBER)
            tip_top += Inches(0.18)
            for t in tippings[:2]:
                cond = t.get("condition", "")[:80]
                if len(t.get("condition", "")) > 80:
                    cond += "…"
                h_days = t.get("horizon_days", 90)
                _text(slide, f"Si {cond} (J+{h_days})",
                      bx + Inches(0.12), tip_top,
                      block_w - Inches(0.24), Inches(0.18),
                      size=8, color=_rgb("555555"), italic=True)
                tip_top += Inches(0.18)
                if tip_top >= block_top + block_h - Inches(0.05):
                    break


# ─── SLIDE 8 : EXÉCUTION (ROADMAP 30/60/90) ──────────────────────────────────

def _slide_execution(prs, edm, result: dict, company: str, date_str: str, page: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header_band(slide, "EXÉCUTION", company)
    _footer_band(slide, page, company)
    _slide_title(slide, "Roadmap d'exécution — 30 / 60 / 90 jours")

    phases_edm = edm.roadmap_90_days or []
    plan_items = result.get("plan_action_30_60_90") or []

    def _phase_items(horizon: str) -> List[str]:
        for phase in phases_edm:
            if str(phase.horizon) == horizon:
                # C3 fix: appliquer _sm() pour nettoyer les ** markdown LLM
                return [_sm(a.decision) for a in phase.actions[:5]]
        if plan_items:
            items = []
            for p in plan_items:
                h = str(getattr(p, "horizon", None) or (p.get("horizon", "") if isinstance(p, dict) else ""))
                if h == horizon:
                    act = (getattr(p, "action", None) or p.get("action", "")) if isinstance(p, dict) else getattr(p, "action", "")
                    if act:
                        # C3 fix: strip markdown ici aussi
                        items.append(_sm(str(act)))
            return items[:5]
        return []

    cols = [("30 JOURS", "30", BLUE), ("60 JOURS", "60", AMBER), ("90 JOURS", "90", GREEN)]
    col_w = CW / 3 - Inches(0.1)
    cy = MT + Inches(0.55)

    for i, (label, horizon, c) in enumerate(cols):
        lx = ML + i * (col_w + Inches(0.15))
        _rect(slide, lx, cy, col_w, Inches(0.45), fill_color=c)
        _text(slide, label, lx, cy + Pt(4), col_w, Inches(0.4),
              size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        items = _phase_items(horizon) or ["—"]
        item_y = cy + Inches(0.55)
        for txt in items:
            item_h = Inches(0.65)
            _rect(slide, lx, item_y, col_w, item_h - Pt(2), fill_color=_rgb("F5F8FF"), line_color=LGRAY)
            _rect(slide, lx + Inches(0.12), item_y + Inches(0.22), Inches(0.08), Inches(0.08), fill_color=c)
            _text(slide, txt, lx + Inches(0.28), item_y + Pt(4),
                  col_w - Inches(0.38), item_h - Pt(8), size=12, color=DARK)
            item_y += item_h + Pt(3)
            if item_y > SH - Inches(0.65):
                break


# ─── SLIDE 9 : SIMULATION (Action vs Inaction) ────────────────────────────────

def _slide_simulation(prs, edm, result: dict, company: str, date_str: str, page: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header_band(slide, "SIMULATION", company)
    _footer_band(slide, page, company)
    _slide_title(slide, "Impact action vs inaction", "Trajectoire financière sur 12 mois")

    series_a = edm.action_series or []
    series_b = edm.do_nothing_series or []
    has_data = (len(series_a) == 12 and len(series_b) == 12
                and any(v != 0 for v in series_a))

    # ── Scenario cards (fixed at bottom) ────────────────────────────────────
    scenarios = result.get("scenarios") or edm.scenarios or []
    scen_map = {}
    for sc in scenarios:
        if isinstance(sc, dict):
            nom, lbl, desc = sc.get("nom", ""), sc.get("label", ""), sc.get("description", "")
        else:
            nom = getattr(sc, "nom", "")
            lbl = getattr(sc, "label", "")
            desc = getattr(sc, "description", "")
        n = nom.lower()
        if "best" in n or "meilleur" in n:
            scen_map["best"] = (lbl, desc)
        elif "likely" in n or "probable" in n or "most" in n:
            scen_map["likely"] = (lbl, desc)
        elif "worst" in n or "pire" in n:
            scen_map["worst"] = (lbl, desc)

    scen_data = [
        ("Meilleur cas",         scen_map.get("best",   ("", ""))[1], GREEN),
        ("Cas le plus probable", scen_map.get("likely", ("", ""))[1], BLUE),
        ("Pire cas",             scen_map.get("worst",  ("", ""))[1], RED),
    ]
    SCEN_Y    = int(MT + Inches(3.65))
    col_h_sc  = int(SH - MB - SCEN_Y - Inches(0.05))
    sw_col    = int(CW / 3 - Inches(0.12))
    for i, (slbl, desc, c) in enumerate(scen_data):
        lx = int(ML + i * (sw_col + Inches(0.18)))
        _rect(slide, lx, SCEN_Y, sw_col, col_h_sc, fill_color=_rgb("F5F8FF"), line_color=c)
        _rect(slide, lx, SCEN_Y, sw_col, int(Inches(0.045)), fill_color=c)
        _text(slide, slbl.upper(), lx + int(Inches(0.12)), SCEN_Y + int(Inches(0.07)),
              sw_col - int(Inches(0.24)), int(Inches(0.28)), size=9, bold=True, color=c)
        _text(slide, desc if desc else "—",
              lx + int(Inches(0.12)), SCEN_Y + int(Inches(0.41)),
              sw_col - int(Inches(0.24)), col_h_sc - int(Inches(0.5)),
              size=10, color=DARK)

    if not has_data:
        _text(slide, "Données de simulation non disponibles",
              int(ML), int(MT + Inches(0.8)), int(CW), int(Inches(0.8)), size=16, color=GRAY)
        return

    # ── Line chart ──────────────────────────────────────────────────────────
    # Layout: y-axis labels on left, chart area fills the rest
    YLAB_W = int(Inches(0.58))
    CX     = int(ML + YLAB_W)               # chart left x
    CY     = int(MT + Inches(0.80))         # chart top y
    CW_C   = int(CW - YLAB_W)              # chart width
    CB     = int(MT + Inches(3.38))         # chart bottom y
    CH_C   = CB - CY                        # chart height
    XPAD   = int(Inches(0.22))             # horizontal inner padding
    PLOT_X = CX + XPAD                     # plot area left
    PLOT_W = CW_C - 2 * XPAD              # plot area width

    # Value range with 12% headroom top & bottom
    all_vals = list(series_a) + list(series_b) + [0]
    v_min0, v_max0 = min(all_vals), max(all_vals)
    v_span0 = v_max0 - v_min0 or 1
    v_min  = v_min0 - v_span0 * 0.12
    v_max  = v_max0 + v_span0 * 0.12
    v_span = v_max - v_min

    def _vy(val: float) -> int:
        """Data value → EMU y-coordinate (high value = low y)"""
        return int(CY + CH_C * (1.0 - (val - v_min) / v_span))

    def _vx(mi: int) -> int:
        """Month index 0-11 → EMU x-coordinate"""
        return int(PLOT_X + mi / 11.0 * PLOT_W)

    # Chart background
    _rect(slide, CX, CY, CW_C, CH_C, fill_color=_rgb("F8FAFC"), line_color=_rgb("DDE2EA"))

    # Horizontal grid lines at 5 levels + labels
    for pct in [0.0, 0.25, 0.5, 0.75, 1.0]:
        gval = v_min + pct * v_span
        gy   = _vy(gval)
        if CY <= gy <= CB:
            is_zero  = abs(gval) < v_span * 0.03
            gcol     = _rgb("9AA5B4") if is_zero else _rgb("E0E5EC")
            g_thick  = int(Pt(1.2)) if is_zero else int(Pt(0.6))
            _rect(slide, CX, gy, CW_C, max(1, g_thick), fill_color=gcol)
            # Y-axis label — no € suffix to avoid wrap in narrow 0.52" box
            _text(slide, _fmt_yax(gval),
                  int(ML), gy - int(Inches(0.14)),
                  YLAB_W - int(Inches(0.06)), int(Inches(0.28)),
                  size=8, color=_rgb("7A8699"), align=PP_ALIGN.RIGHT)

    # Month tick marks on x-axis
    for mi in range(12):
        tx = _vx(mi)
        _rect(slide, tx, CB, int(Pt(0.75)), int(Inches(0.06)), fill_color=_rgb("9AA5B4"))

    # Draw do_nothing series (red, behind)
    for i in range(11):
        _draw_line_segment(slide,
                           _vx(i),   _vy(series_b[i]),
                           _vx(i+1), _vy(series_b[i+1]),
                           RED, int(Pt(2.5)))

    # Draw action series (green, in front)
    for i in range(11):
        _draw_line_segment(slide,
                           _vx(i),   _vy(series_a[i]),
                           _vx(i+1), _vy(series_a[i+1]),
                           GREEN, int(Pt(3.0)))

    # Dots + callout values at M3, M6, M9, M12
    DOT = int(Inches(0.075))
    for mi in [2, 5, 8, 11]:
        # Green dot (action)
        xa, ya = _vx(mi), _vy(series_a[mi])
        _rect(slide, xa - DOT // 2, ya - DOT // 2, DOT, DOT, fill_color=GREEN)
        _text(slide, _fmt_chart(series_a[mi]),
              xa - int(Inches(0.28)), ya - int(Inches(0.32)),
              int(Inches(0.56)), int(Inches(0.26)),
              size=8, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
        # Red dot (do nothing)
        xb, yb = _vx(mi), _vy(series_b[mi])
        _rect(slide, xb - DOT // 2, yb - DOT // 2, DOT, DOT, fill_color=RED)
        _text(slide, _fmt_chart(series_b[mi]),
              xb - int(Inches(0.28)), yb + int(Inches(0.07)),
              int(Inches(0.56)), int(Inches(0.26)),
              size=8, bold=True, color=RED, align=PP_ALIGN.CENTER)

    # X-axis labels — 0.48" box, centred on tick mark
    XLAB_Y = CB + int(Inches(0.05))
    for mi, lbl in [(0, "M1"), (2, "M3"), (5, "M6"), (8, "M9"), (11, "M12")]:
        _text(slide, lbl,
              _vx(mi) - int(Inches(0.24)), XLAB_Y,
              int(Inches(0.48)), int(Inches(0.22)),
              size=8, color=_rgb("7A8699"), align=PP_ALIGN.CENTER)

    # Legend box (top-right of chart)
    LEG_X = CX + CW_C - int(Inches(2.08))
    LEG_Y = CY + int(Inches(0.1))
    LEG_W = int(Inches(1.98))
    LEG_H = int(Inches(0.62))
    _rect(slide, LEG_X, LEG_Y, LEG_W, LEG_H,
          fill_color=WHITE, line_color=_rgb("CCCCCC"))
    _draw_line_segment(slide,
                       LEG_X + int(Inches(0.1)),  LEG_Y + int(Inches(0.19)),
                       LEG_X + int(Inches(0.38)), LEG_Y + int(Inches(0.19)),
                       GREEN, int(Pt(3)))
    _text(slide, "Avec décisions",
          LEG_X + int(Inches(0.44)), LEG_Y + int(Inches(0.07)),
          int(Inches(1.48)), int(Inches(0.24)), size=9, color=GREEN)
    _draw_line_segment(slide,
                       LEG_X + int(Inches(0.1)),  LEG_Y + int(Inches(0.44)),
                       LEG_X + int(Inches(0.38)), LEG_Y + int(Inches(0.44)),
                       RED, int(Pt(2.5)))
    _text(slide, "Sans action",
          LEG_X + int(Inches(0.44)), LEG_Y + int(Inches(0.33)),
          int(Inches(1.48)), int(Inches(0.24)), size=9, color=RED)


# ─── SLIDE 10 : PROJECTION ────────────────────────────────────────────────────

def _slide_projection(prs, edm, result: dict, company: str, date_str: str, page: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header_band(slide, "PROJECTION", company)
    _footer_band(slide, page, company)
    _slide_title(slide, "Trajectoire financière — 12 mois",
                 "Si les décisions prioritaires sont engagées maintenant")

    series = edm.monthly_projection or []
    has_data = len(series) == 12 and any(v != 0 for v in series)

    # ── Layout — subtitle bottom = MT + 0.65" → chart starts at MT + 0.78" ──
    CONTENT_Y = MT + Inches(0.78)   # safe margin below subtitle
    CHART_H   = Inches(2.55)
    CHART_BOT = CONTENT_Y + CHART_H

    break_even = None

    if has_data:
        min_v = min(series)
        max_v = max(series)
        # 8% vertical padding for breathing room
        pad = (max_v - min_v) * 0.08 if max_v != min_v else abs(max_v) * 0.08 or 1
        plot_min = (min_v - pad) if min_v < 0 else 0
        plot_max = max_v + pad
        plot_range = plot_max - plot_min or 1

        # Zero baseline
        y_zero = int(CONTENT_Y + CHART_H * (plot_max / plot_range))

        # Thin baseline rule
        bl = slide.shapes.add_shape(1, int(ML), y_zero - int(Pt(1)), int(CW), int(Pt(2)))
        bl.fill.solid()
        bl.fill.fore_color.rgb = _rgb("9AACBA")
        bl.line.fill.background()

        # Y-axis scale labels + subtle horizontal gridlines
        for pct in [0.0, 0.5, 1.0]:
            gval = plot_min + pct * plot_range
            gy = int(CONTENT_Y + CHART_H * (1.0 - (gval - plot_min) / plot_range))
            if int(CONTENT_Y) <= gy <= int(CHART_BOT):
                # Faint gridline (drawn before bars so bars appear on top)
                _rect(slide, int(ML), gy, int(CW), int(Pt(0.5)),
                      fill_color=_rgb("DEE3EA"))
                # Y label — no € to avoid wrap in narrow left margin
                _text(slide, _fmt_yax(gval),
                      0, gy - int(Inches(0.16)),
                      int(ML) - int(Inches(0.04)), int(Inches(0.28)),
                      size=7, color=_rgb("9AA5B4"), align=PP_ALIGN.RIGHT)

        slot_w = int(CW / 12)
        bar_w  = int(slot_w * 0.66)
        bar_off = (slot_w - bar_w) // 2

        for i, v in enumerate(series):
            fc  = GREEN if v >= 0 else RED
            bx  = int(ML) + i * slot_w + bar_off
            bar_h_emu = max(int(Pt(3)), int(CHART_H * abs(v) / plot_range))
            by  = (y_zero - bar_h_emu) if v >= 0 else y_zero
            _rect(slide, bx, by, bar_w, bar_h_emu, fill_color=fc)

            # M1–M12 label
            _text(slide, f"M{i+1}",
                  int(ML) + i * slot_w, int(CHART_BOT) + int(Inches(0.04)),
                  slot_w, int(Inches(0.22)),
                  size=9, color=GRAY, align=PP_ALIGN.CENTER)

            # Compact value label
            abs_v = abs(v)
            pfx   = "−" if v < 0 else ""   # proper minus sign
            if abs_v >= 1_000_000:
                lbl = f"{pfx}{abs_v/1_000_000:.1f}M€".replace(".", ",")
            elif abs_v >= 1_000:
                lbl = f"{pfx}{int(abs_v/1_000)}K€"
            else:
                lbl = f"{pfx}{int(abs_v)}€"
            _text(slide, lbl,
                  int(ML) + i * slot_w, int(CHART_BOT) + int(Inches(0.27)),
                  slot_w, int(Inches(0.22)),
                  size=9, bold=True, color=fc, align=PP_ALIGN.CENTER)

        break_even = next((i + 1 for i, v in enumerate(series) if v >= 0), None)
        msg = (f"Retour à l’équilibre estimé au mois {break_even}, "
               f"sous réserve d’engagement des décisions."
               if break_even
               else "Le retour à l’équilibre n’est pas atteint sur 12 mois.")
        note_y = CHART_BOT + Inches(0.56)
        _text(slide, msg, ML, note_y, int(CW * 0.58), Inches(0.32),
              size=12, color=GRAY, italic=True)
    else:
        _text(slide, "—", ML, CONTENT_Y, CW, Inches(1), size=22, color=GRAY)
        note_y = CONTENT_Y + Inches(1.1)

    # ── Impact décisions box — bottom right ───────────────────────────────────
    # mn1 fix: utiliser decisions_net (impact des décisions exécutées) plutôt
    # que coi.per_year (COI total) pour cohérence avec la courbe du graphique
    decisions_net = sum((d.annual_impact or 0) for d in (edm.executive_decisions or []))
    if decisions_net:
        box_x = ML + int(CW * 0.62)
        box_w = int(CW * 0.38)
        box_y = CHART_BOT + Inches(0.51) if has_data else CONTENT_Y + Inches(1.05)
        box_h = Inches(1.08)
        box_color = GREEN if decisions_net > 0 else RED
        _rect(slide, box_x, box_y, box_w, box_h,
              fill_color=_rgb("F0FBF4"), line_color=box_color)
        _text(slide, "Impact si vous agissez · 12 mois",
              box_x + Inches(0.12), box_y + Inches(0.1),
              box_w - Inches(0.24), Inches(0.3),
              size=10, bold=True, color=box_color, align=PP_ALIGN.CENTER)
        _text(slide, _fmt_auto(abs(decisions_net)),
              box_x + Inches(0.1), box_y + Inches(0.38),
              box_w - Inches(0.2), Inches(0.58),
              size=30, bold=True, color=box_color, align=PP_ALIGN.CENTER)


# ─── SHARED WATERFALL RENDERER ───────────────────────────────────────────────

def _draw_waterfall(slide, items, *,
                    chart_x, chart_w, chart_top, chart_bot,
                    ylabel_w, ml):
    """
    items : list of (lbl1, lbl2, delta, kind)
      kind : "anchor" | "positive" | "negative"
    Draws bars, connectors, y-labels, x-labels.
    Returns nothing (mutates slide in place).
    """
    N      = len(items)
    CH     = chart_bot - chart_top

    bar_w  = int(chart_w / (N + N * 0.70))
    gap    = int(bar_w * 0.70)
    side_m = int((chart_w - N * bar_w - (N - 1) * gap) / 2)

    def bx(i):
        return chart_x + side_m + i * (bar_w + gap)

    # Y scale — include all running totals
    levels = [0.0]
    r = 0.0
    for lbl1, lbl2, delta, kind in items:
        if kind == "anchor":
            levels.append(delta)
            r = delta
        else:
            r += delta
            levels.append(r)

    y_min_v = min(levels)
    y_max_v = max(levels)
    pad_lo  = abs(y_min_v) * 0.28 if y_min_v != 0 else abs(y_max_v) * 0.10
    pad_hi  = abs(y_max_v) * 0.22 if y_max_v != 0 else abs(y_min_v) * 0.10
    y_min_v -= pad_lo
    y_max_v += pad_hi
    y_min_v  = min(y_min_v, -30_000)
    y_max_v  = max(y_max_v,  30_000)
    Y_SPAN   = y_max_v - y_min_v

    def vy(val):
        return int(chart_bot - (val - y_min_v) / Y_SPAN * CH)

    y_zero = vy(0)

    NAVY_ANC = _rgb("0D2B6E")
    GRN_BAR  = _rgb("1B6B3A")
    RED_BAR  = _rgb("B71C1C")
    GRID_CLR = _rgb("DDE3EC")
    ZERO_CLR = _rgb("5A6A7E")
    CONN_CLR = _rgb("8A9DBB")

    # Grid lines + Y labels
    n_grids = 5
    for gi in range(n_grids):
        gval = y_min_v + gi * Y_SPAN / (n_grids - 1)
        gy   = vy(gval)
        _rect(slide, chart_x, gy, chart_w, max(1, int(Inches(0.010))), fill_color=GRID_CLR)
        _text(slide, _fmt_yax(gval),
              int(ml), gy - int(Inches(0.14)),
              ylabel_w - int(Inches(0.10)), int(Inches(0.28)),
              size=8, color=_rgb("7A8699"), align=PP_ALIGN.RIGHT)

    # Zero line
    _rect(slide, chart_x, y_zero, chart_w, max(2, int(Inches(0.018))), fill_color=ZERO_CLR)
    _text(slide, "0",
          int(ml), y_zero - int(Inches(0.14)),
          ylabel_w - int(Inches(0.10)), int(Inches(0.28)),
          size=8, bold=True, color=ZERO_CLR, align=PP_ALIGN.RIGHT)

    # Bars + connectors
    running     = 0.0
    prev_right  = None
    prev_run_y  = None

    for i, (lbl1, lbl2, delta, kind) in enumerate(items):
        x      = bx(i)
        next_x = bx(i + 1) if i + 1 < N else None

        if kind == "anchor":
            y_val  = vy(delta)
            y_z    = vy(0)
            b_top  = min(y_val, y_z)
            b_h    = max(abs(y_val - y_z), int(Inches(0.05)))
            _rect(slide, x, b_top, bar_w, b_h, fill_color=NAVY_ANC)

            lbl_y = (y_val - int(Inches(0.28))) if delta >= 0 else (y_z + int(Inches(0.06)))
            _text(slide, _fmt_chart(delta),
                  x - int(Inches(0.10)), lbl_y,
                  bar_w + int(Inches(0.20)), int(Inches(0.24)),
                  size=10, bold=True, color=NAVY_ANC, align=PP_ALIGN.CENTER)

            # Connector from previous running level
            if prev_right is not None and prev_run_y is not None:
                conn_w = x - prev_right
                if conn_w > 0:
                    _rect(slide, prev_right, prev_run_y, conn_w,
                          max(1, int(Inches(0.012))), fill_color=CONN_CLR)

            running     = delta
            prev_right  = x + bar_w
            prev_run_y  = vy(delta)

        else:
            base_y = vy(running)
            top_y  = vy(running + delta)
            b_top  = min(base_y, top_y)
            b_h    = max(abs(base_y - top_y), int(Inches(0.05)))
            clr    = GRN_BAR if delta >= 0 else RED_BAR
            _rect(slide, x, b_top, bar_w, b_h, fill_color=clr)

            lbl_y = (top_y - int(Inches(0.28))) if delta >= 0 else (base_y + int(Inches(0.06)))
            _text(slide, _fmt_chart(delta),
                  x - int(Inches(0.10)), lbl_y,
                  bar_w + int(Inches(0.20)), int(Inches(0.24)),
                  size=9, bold=True, color=clr, align=PP_ALIGN.CENTER)

            # Connector from previous running level
            if prev_right is not None and prev_run_y is not None:
                conn_w = x - prev_right
                if conn_w > 0:
                    _rect(slide, prev_right, prev_run_y, conn_w,
                          max(1, int(Inches(0.012))), fill_color=CONN_CLR)

            running    += delta
            prev_right  = x + bar_w
            prev_run_y  = vy(running)

        # X-axis label
        lbl_base_y = chart_bot + int(Inches(0.10))
        _text(slide, lbl1,
              x - int(Inches(0.20)), lbl_base_y,
              bar_w + int(Inches(0.40)), int(Inches(0.24)),
              size=8.5, bold=(kind == "anchor"), color=DARK, align=PP_ALIGN.CENTER)
        if lbl2:
            _text(slide, lbl2,
                  x - int(Inches(0.20)), lbl_base_y + int(Inches(0.24)),
                  bar_w + int(Inches(0.40)), int(Inches(0.20)),
                  size=7.5, color=_rgb("5A6A7E"), align=PP_ALIGN.CENTER)


# ─── SLIDE 10a : BRIDGE HISTORIQUE ───────────────────────────────────────────

def _slide_bridge_historique(prs, edm, result: dict, company: str, date_str: str, page: int):
    """
    Bridge historique : EBITDA Normatif → value_destroyers (barres rouges) → EBITDA Actuel.
    Répond à "d'où vient la perte de valeur sur la période analysée ?"
    """
    import re as _re, datetime as _dt

    # ── Année + mois de référence ───────────────────────────────────────────────
    _ym = _re.search(r"\b(20\d{2})\b", date_str or "")
    year_n = int(_ym.group(1)) if _ym else _dt.datetime.now().year
    _MOIS_ABBR = {
        "janvier": "Jan.", "février": "Fév.", "mars": "Mar.",
        "avril": "Avr.", "mai": "Mai", "juin": "Juin",
        "juillet": "Juil.", "août": "Août", "septembre": "Sep.",
        "octobre": "Oct.", "novembre": "Nov.", "décembre": "Déc.",
    }
    month_abbr = next((abbr for m, abbr in _MOIS_ABBR.items()
                       if m in (date_str or "").lower()), str(year_n))

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _footer_band(slide, page, company)

    # ── Data ──────────────────────────────────────────────────────────────────
    def _parse_kpi(kpi_list, keyword):
        for k in (kpi_list or []):
            if keyword.lower() in (k.get("label") or "").lower():
                s = (k.get("value") or "").replace("−", "-").replace("−", "-")
                s = s.replace(" ", "").replace("\xa0", "").replace(" ", "").replace(",", ".")
                m = _re.search(r"([+-]?\d+\.?\d*)([KkMm])?", s)
                if m:
                    n = float(m.group(1))
                    mult = (m.group(2) or "").upper()
                    return n * (1_000 if mult == "K" else 1_000_000 if mult == "M" else 1)
        return None

    kpi_dash      = result.get("ceo_dashboard") or []
    ebitda_actuel = _parse_kpi(kpi_dash, "ebitda") or 0.0
    destroyers    = edm.value_destroyers or []
    total_destr   = sum(vd.annual_impact or 0 for vd in destroyers)
    ebitda_norm   = ebitda_actuel - total_destr

    # ── Logique dynamique : entreprise dégradée vs saine ──────────────────────
    is_degraded = total_destr < 0
    if is_degraded:
        header_txt  = "ANALYSE DE LA VALEUR PERDUE"
        slide_ttl   = "D'où vient la perte de valeur ?"
        callout_txt = f"Valeur détruite : {_fmt_chart(total_destr)} / an"
        callout_clr = _rgb("B71C1C")
    else:
        header_txt  = "ANALYSE DE LA CRÉATION DE VALEUR"
        slide_ttl   = "D'où vient la création de valeur ?"
        callout_txt = f"Valeur créée : {_fmt_chart(total_destr)} / an"
        callout_clr = _rgb("1B6B3A")

    _header_band(slide, header_txt, company)

    # ── Title ─────────────────────────────────────────────────────────────────
    # mn4 fix: _fmt_chart arrondissait 2,165M et 2,242M tous les deux à "+2,2M€"
    # → subtitle affichait "normatif (+2,2M€) à actuel (+2,2M€)" (valeurs identiques).
    # Forcer le format K€ pour que les deux bornes soient distinguables.
    def _fmt_ke(v: float) -> str:
        sign = "+" if v >= 0 else ""
        return f"{sign}{int(round(v / 1_000)):,}K€".replace(",", " ")
    _slide_title(slide, slide_ttl,
                 f"Exercice {year_n} — de l'EBITDA normatif ({_fmt_ke(ebitda_norm)}) "
                 f"à la situation actuelle ({_fmt_ke(ebitda_actuel)})")

    # ── Items ─────────────────────────────────────────────────────────────────
    def _split_label(txt, n=18):
        """Découpe un label sur 2 lignes (lbl1, lbl2) au lieu de tronquer avec '…'."""
        txt = _re.split(r"[—\(]", txt)[0].strip()
        if len(txt) <= n:
            return txt, ""
        split_pos = txt[:n].rfind(" ")
        if split_pos < 6:
            split_pos = n
        line1 = txt[:split_pos].rstrip()
        line2 = txt[split_pos:].strip()
        if len(line2) > n:
            line2 = line2[:n - 1].rstrip() + "…"
        return line1, line2

    items = [("EBITDA", f"Normatif · {month_abbr} {year_n}", ebitda_norm, "anchor")]
    for vd in destroyers[:5]:
        delta = vd.annual_impact or 0
        lbl1, lbl2 = _split_label(_sm(vd.name))
        items.append((lbl1, lbl2, delta, "negative" if delta < 0 else "positive"))
    items.append(("EBITDA", f"Actuel · {month_abbr} {year_n}", ebitda_actuel, "anchor"))

    # ── Layout + rendu ────────────────────────────────────────────────────────
    YLABEL_W = int(Inches(1.20))
    CX       = int(ML + YLABEL_W)
    CHART_W  = int(SW - ML - MR - YLABEL_W)
    CT       = int(MT + Inches(1.30))
    CB       = int(SH - MB - Inches(0.55))

    _draw_waterfall(slide, items,
                    chart_x=CX, chart_w=CHART_W,
                    chart_top=CT, chart_bot=CB,
                    ylabel_w=YLABEL_W, ml=ML)

    # ── Callout dynamique (top-right) ─────────────────────────────────────────
    _text(slide, callout_txt,
          CX + CHART_W - int(Inches(3.2)), CT - int(Inches(0.02)),
          int(Inches(3.10)), int(Inches(0.28)),
          size=9, bold=True, color=callout_clr, align=PP_ALIGN.RIGHT)


# ─── SLIDE 10b : BRIDGE FINANCIER ────────────────────────────────────────────

def _slide_bridge_financier(prs, edm, result: dict, company: str, date_str: str, page: int):
    """
    Waterfall / bridge chart  : EBITDA N  →  décisions  →  EBITDA cible N+1.
    Chaque décision apparaît comme une barre verte flottante.
    Barres d'ancrage (départ / arrivée) en bleu marine.
    """
    import re as _re, datetime as _dt

    # ── Années + mois de référence ────────────────────────────────────────────
    _ym = _re.search(r"\b(20\d{2})\b", date_str or "")
    year_n  = int(_ym.group(1)) if _ym else _dt.datetime.now().year
    year_n1 = year_n + 1
    _MOIS_ABBR_BF = {
        "janvier": "Jan.", "février": "Fév.", "mars": "Mar.",
        "avril": "Avr.", "mai": "Mai", "juin": "Juin",
        "juillet": "Juil.", "août": "Août", "septembre": "Sep.",
        "octobre": "Oct.", "novembre": "Nov.", "décembre": "Déc.",
    }
    month_abbr = next((abbr for m, abbr in _MOIS_ABBR_BF.items()
                       if m in (date_str or "").lower()), str(year_n))

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header_band(slide, "BRIDGE FINANCIER", company)
    _footer_band(slide, page, company)

    # ── Extraire EBITDA depuis le KPI dashboard ─────────────────────────────
    def _parse_kpi(kpi_list, keyword):
        for k in (kpi_list or []):
            if keyword.lower() in (k.get("label") or "").lower():
                s = (k.get("value") or "").replace("−", "-").replace("−", "-")
                s = s.replace(" ", "").replace("\xa0", "").replace(" ", "").replace(",", ".")
                m = _re.search(r"([+-]?\d+\.?\d*)([KkMm])?", s)
                if m:
                    n = float(m.group(1))
                    mult = (m.group(2) or "").upper()
                    return n * (1_000 if mult == "K" else 1_000_000 if mult == "M" else 1)
        return None

    kpi_dash = result.get("ceo_dashboard") or []
    ebitda_start = _parse_kpi(kpi_dash, "ebitda")
    if ebitda_start is None:
        # Fallback : utiliser le COI mensuel comme proxy de la perte annuelle
        coi = edm.cost_of_inaction
        ebitda_start = float(getattr(coi, "per_year", 0) or 0) / 12

    decisions = edm.executive_decisions or []
    total_impact = sum(d.annual_impact or 0 for d in decisions)
    ebitda_end = ebitda_start + total_impact

    # ── Titre slide ──────────────────────────────────────────────────────────
    n_dec = len(decisions[:5])
    subtitle = (
        f"FY {year_n} → FY {year_n1} — impact annuel des {n_dec} décisions prioritaires — "
        f"{_fmt_chart(ebitda_start)} → {_fmt_chart(ebitda_end)} d'EBITDA"
    )
    _slide_title(slide, "De la situation actuelle à l'objectif", subtitle)

    # ── Bridge items : [(label_line1, label_line2, delta, kind)] ─────────────
    # kind : "anchor" | "positive" | "negative"
    def _split_label_bf(txt, n=18):
        """Découpe un label sur 2 lignes (lbl1, lbl2) au lieu de tronquer avec '…'."""
        txt = _re.split(r"[—\(,]", txt)[0].strip()
        if len(txt) <= n:
            return txt, ""
        split_pos = txt[:n].rfind(" ")
        if split_pos < 6:
            split_pos = n
        line1 = txt[:split_pos].rstrip()
        line2 = txt[split_pos:].strip()
        if len(line2) > n:
            line2 = line2[:n - 1].rstrip() + "…"
        return line1, line2

    items = [("EBITDA", f"Actuel · {month_abbr} {year_n}", ebitda_start, "anchor")]
    for d in decisions[:5]:
        delta = d.annual_impact or 0
        kind  = "positive" if delta >= 0 else "negative"
        lbl1, lbl2 = _split_label_bf(_sm(d.decision))
        items.append((lbl1, lbl2, delta, kind))
    items.append(("EBITDA", f"Cible · {month_abbr} {year_n1}", ebitda_end, "anchor"))

    # ── Layout + shared renderer ──────────────────────────────────────────────
    YLABEL_W = int(Inches(1.20))
    CX       = int(ML + YLABEL_W)
    CHART_W  = int(SW - ML - MR - YLABEL_W)
    CT       = int(MT + Inches(1.30))
    CB       = int(SH - MB - Inches(0.55))

    _draw_waterfall(slide, items,
                    chart_x=CX, chart_w=CHART_W,
                    chart_top=CT, chart_bot=CB,
                    ylabel_w=YLABEL_W, ml=ML)

    # ── Total-impact callout ──────────────────────────────────────────────────
    _text(slide, f"Impact total : {_fmt_chart(total_impact)} / an",
          CX + CHART_W - int(Inches(3.0)), CT - int(Inches(0.02)),
          int(Inches(2.90)), int(Inches(0.28)),
          size=9, bold=True, color=_rgb("1B6B3A"), align=PP_ALIGN.RIGHT)


# ─── SLIDE 11 : CRÉATION DE VALEUR ───────────────────────────────────────────

def _slide_creation_valeur(prs, edm, result: dict, company: str, date_str: str, page: int):
    """
    Slide : Potentiel de création de valeur.
    Phase 1 — Récupération structurelle (décisions déjà identifiées).
    Phase 2 — Leviers de croissance offensive (à activer).
    Adaptatif : formulation crise (redressement → croissance) ou croissance (accélération).
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header_band(slide, "CRÉATION DE VALEUR", company)
    _footer_band(slide, page, company)

    is_growth = _is_growth_context(edm, result)

    if is_growth:
        title    = "Leviers de création de valeur"
        subtitle = "Optimisation et accélération de la performance"
        p1_label = "Optimisations opérationnelles"
        p2_label = "Leviers de croissance offensive"
        p2_timeline = "6–18 mois"
    else:
        title    = "Potentiel de création de valeur"
        subtitle = "Du redressement structurel à la croissance durable"
        p1_label = "Récupération structurelle"
        p2_label = "Accélération offensive"
        p2_timeline = "6–18 mois"

    _slide_title(slide, title, subtitle)

    # ── Data ────────────────────────────────────────────────────────────────────
    decisions  = edm.executive_decisions or []
    p1_total   = sum((d.annual_impact or 0) for d in decisions)
    p1_items   = [(_sm(d.decision)[:62] + ("…" if len(_sm(d.decision)) > 62 else ""),
                   _fmt_auto(d.annual_impact)) for d in decisions[:4]]

    # Phase 2 : leviers de croissance (strategic_levers ou synthèse contextuelle)
    raw_levers = edm.strategic_levers or []
    if raw_levers:
        levers_data = [(getattr(l, "name", "—"),
                        getattr(l, "description", ""),
                        _fmt_auto(getattr(l, "estimated_impact", None)))
                       for l in raw_levers[:3]]
        p2_low  = int(sum((getattr(l, "estimated_impact", 0) or 0) for l in raw_levers) * 0.70)
        p2_high = int(sum((getattr(l, "estimated_impact", 0) or 0) for l in raw_levers) * 1.20)
    elif is_growth:
        levers_data = [
            ("Pricing power & mix produit",
             "Augmenter le panier moyen de 10–15 % sur les segments premium",
             "+200K€–500K€"),
            ("Expansion commerciale",
             "Nouveaux segments clients ou zones géographiques activables",
             "+300K€–800K€"),
            ("Levier digital & automation",
             "Réduire les coûts opérationnels de 8–12 % via l'automatisation",
             "+150K€–400K€"),
        ]
        p2_low, p2_high = 650_000, 1_700_000
    else:
        levers_data = [
            ("Rétablir la marge brute cible",
             "Revenir à 34–36 % via optimisation du mix et renégociation tarifaire",
             "+300K€–450K€"),
            ("Plan de conquête commerciale",
             "Réactiver la croissance du CA (+15 %) avec la nouvelle grille commerciale",
             "+400K€–600K€"),
            ("Réinvestir la trésorerie libérée",
             "Financer R&D ou acquisitions clients grâce au BFR optimisé",
             "+200K€–350K€"),
        ]
        p2_low, p2_high = 900_000, 1_400_000

    vcs = (getattr(edm, "value_creation_statement", None)
           or result.get("value_creation_statement", ""))

    # ── Layout constants ────────────────────────────────────────────────────────
    TOP_Y   = int(MT + Inches(0.90))   # clear title + subtitle
    CARD_H  = int(Inches(2.42))
    GAP     = int(Inches(0.18))
    CARD_W1 = int(CW * 0.46)   # Phase 1 (left)
    CARD_W2 = int(CW * 0.50)   # Phase 2 (right)
    CARD_X2 = int(ML + CARD_W1 + GAP)
    LEV_Y   = TOP_Y + CARD_H + int(Inches(0.14))
    LEV_H   = int(Inches(1.30))
    lev_w   = int((CW - 2 * GAP) / 3)

    AMBER   = _rgb("C47B1A")
    GREEN2  = _rgb("1B6B3A")   # dark green for accent
    TEAL    = _rgb("028090")

    # ── Phase 1 card ────────────────────────────────────────────────────────────
    _rect(slide, int(ML), TOP_Y, CARD_W1, CARD_H,
          fill_color=_rgb("FDFAF2"), line_color=AMBER)
    _rect(slide, int(ML), TOP_Y, CARD_W1, int(Inches(0.045)), fill_color=AMBER)

    _text(slide, p1_label.upper(),
          int(ML) + int(Inches(0.14)), TOP_Y + int(Inches(0.07)),
          CARD_W1 - int(Inches(0.28)), int(Inches(0.28)),
          size=9, bold=True, color=AMBER)

    # Big number — compact format to avoid wrap
    _text(slide, _fmt_chart(p1_total),
          int(ML) + int(Inches(0.12)), TOP_Y + int(Inches(0.35)),
          CARD_W1 - int(Inches(0.24)), int(Inches(0.54)),
          size=28, bold=True, color=AMBER, align=PP_ALIGN.LEFT)
    _text(slide, "/ an  ·  30–90 jours",
          int(ML) + int(Inches(0.14)), TOP_Y + int(Inches(0.85)),
          CARD_W1 - int(Inches(0.28)), int(Inches(0.22)),
          size=9, color=_rgb("7A5210"), align=PP_ALIGN.LEFT)

    # Decision bullet list
    bull_y = TOP_Y + int(Inches(1.15))
    for dec_txt, dec_imp in p1_items:
        _text(slide, f"• {dec_txt}",
              int(ML) + int(Inches(0.14)), bull_y,
              CARD_W1 - int(Inches(0.28)), int(Inches(0.30)),
              size=8.5, color=DARK)
        _text(slide, dec_imp,
              int(ML) + CARD_W1 - int(Inches(1.0)), bull_y,
              int(Inches(0.92)), int(Inches(0.24)),
              size=8.5, bold=True, color=GREEN, align=PP_ALIGN.RIGHT)
        bull_y += int(Inches(0.32))

    # ── Phase 2 card ────────────────────────────────────────────────────────────
    _rect(slide, CARD_X2, TOP_Y, CARD_W2, CARD_H,
          fill_color=_rgb("F3F9F4"), line_color=GREEN2)
    _rect(slide, CARD_X2, TOP_Y, CARD_W2, int(Inches(0.045)), fill_color=GREEN2)

    _text(slide, p2_label.upper(),
          CARD_X2 + int(Inches(0.14)), TOP_Y + int(Inches(0.07)),
          CARD_W2 - int(Inches(0.28)), int(Inches(0.28)),
          size=9, bold=True, color=GREEN2)

    p2_range = f"{_fmt_chart(p2_low)}–{_fmt_chart(p2_high)}"
    _text(slide, p2_range,
          CARD_X2 + int(Inches(0.12)), TOP_Y + int(Inches(0.35)),
          CARD_W2 - int(Inches(0.24)), int(Inches(0.54)),
          size=24, bold=True, color=GREEN2, align=PP_ALIGN.LEFT)
    _text(slide, f"/ an  ·  {p2_timeline}",
          CARD_X2 + int(Inches(0.14)), TOP_Y + int(Inches(0.85)),
          CARD_W2 - int(Inches(0.28)), int(Inches(0.22)),
          size=9, color=_rgb("145227"), align=PP_ALIGN.LEFT)

    # Lever items in Phase 2 card — name + impact only (descriptions in bottom cards)
    lev_y2 = TOP_Y + int(Inches(1.15))
    for lev_name, _lev_desc, lev_impact in levers_data:
        _text(slide, f"• {lev_name}",
              CARD_X2 + int(Inches(0.14)), lev_y2,
              CARD_W2 - int(Inches(1.5)), int(Inches(0.28)),
              size=8.5, bold=True, color=DARK)
        _text(slide, lev_impact,
              CARD_X2 + CARD_W2 - int(Inches(1.2)), lev_y2,
              int(Inches(1.12)), int(Inches(0.28)),
              size=8.5, color=GREEN2, align=PP_ALIGN.RIGHT)
        lev_y2 += int(Inches(0.36))

    # ── Bottom lever cards ───────────────────────────────────────────────────────
    LEVER_ICONS = ["📈", "🔑", "⚙️"]
    for i, (lev_name, lev_desc, lev_impact) in enumerate(levers_data):
        lx = int(ML) + i * (lev_w + GAP)
        _rect(slide, lx, LEV_Y, lev_w, LEV_H,
              fill_color=_rgb("F8FAFF"), line_color=_rgb("C0CEDB"))
        _rect(slide, lx, LEV_Y, lev_w, int(Inches(0.032)), fill_color=TEAL)
        _text(slide, lev_name.upper(),
              lx + int(Inches(0.12)), LEV_Y + int(Inches(0.06)),
              lev_w - int(Inches(0.24)), int(Inches(0.28)),
              size=8.5, bold=True, color=TEAL)
        _text(slide, lev_impact,
              lx + int(Inches(0.12)), LEV_Y + int(Inches(0.34)),
              lev_w - int(Inches(0.24)), int(Inches(0.26)),
              size=11, bold=True, color=GREEN2)
        if lev_desc:
            _text(slide, lev_desc,
                  lx + int(Inches(0.12)), LEV_Y + int(Inches(0.62)),
                  lev_w - int(Inches(0.24)), int(Inches(0.70)),
                  size=8.5, color=DARK)

    # ── Bottom summary strip ────────────────────────────────────────────────────
    STRIP_Y = LEV_Y + LEV_H + int(Inches(0.10))
    STRIP_H = int(Inches(0.34))
    if STRIP_Y + STRIP_H < int(SH - MB - Inches(0.02)):
        total_low  = int(p1_total + p2_low)
        total_high = int(p1_total + p2_high)
        total_str  = f"{_fmt_auto(total_low)}–{_fmt_auto(total_high)}"
        _rect(slide, int(ML), STRIP_Y, int(CW), STRIP_H, fill_color=NAVY)
        _text(slide,
              f"Valeur totale mobilisable sur 24 mois : {total_str}  ·  "
              f"{vcs}" if vcs else f"Valeur totale mobilisable sur 24 mois : {total_str}",
              int(ML) + int(Inches(0.22)), STRIP_Y + int(Inches(0.05)),
              int(CW) - int(Inches(0.44)), STRIP_H - int(Inches(0.10)),
              size=9, bold=True, color=WHITE)


# ─── SLIDE 12 : RISQUES ───────────────────────────────────────────────────────

def _slide_risques(prs, edm, result: dict, company: str, date_str: str, page: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header_band(slide, "RISQUES", company)
    _footer_band(slide, page, company)
    _slide_title(slide, "Risques majeurs identifiés")

    risques_raw = result.get("problemes_critiques") or result.get("alertes") or []
    destroyers = edm.value_destroyers
    rows = []
    # Map decisions to potential mitigations (first word of action)
    dec_actions = []
    for d in (destroyers or []):
        nm = (d.name or "").strip()
        if nm:
            dec_actions.append(nm)

    def _prob_from_sev(sev: str) -> str:
        s = sev.lower()
        if "élevé" in s or "eleve" in s or "critique" in s or "fort" in s:
            return "Haute"
        if "faible" in s or "faib" in s:
            return "Basse"
        return "Modérée"

    def _mitigation_for(desc: str) -> str:
        """Try to find the best matching decision as mitigation."""
        dl = desc.lower()
        for da in dec_actions:
            if any(w in dl for w in da.lower().split()[:3]):
                return da[:55] + ("…" if len(da) > 55 else "")
        return dec_actions[0][:55] if dec_actions else "Plan d'action à définir"

    # columns: desc, severite, probabilite, horizon, mitigation
    for r in risques_raw[:6]:
        if isinstance(r, dict):
            desc = (r.get("description") or str(r))
            sev  = r.get("severite", "Moyen")
            hor  = r.get("horizon", "Court terme")
        else:
            desc = str(r)
            sev  = "Moyen"
            hor  = "Court terme"
        prob = _prob_from_sev(sev)
        mitig = _mitigation_for(desc)
        rows.append([desc, sev, prob, hor, mitig])
    if not rows and destroyers:
        for d in destroyers[:5]:
            sev  = "Élevé" if d.annual_impact and abs(d.annual_impact) > 500_000 else "Moyen"
            prob = _prob_from_sev(sev)
            mitig = (d.name or "À définir")[:55]
            rows.append([d.name, sev, prob, "Immédiat", mitig])
    if not rows:
        rows = [["—", "—", "—", "—", "—"]]

    # ── Table — 5 columns ─────────────────────────────────────────────────────
    HEADERS = ["Risque identifié", "Sévérité", "Probabilité", "Horizon", "Mitigation prioritaire"]
    col_pct  = [0.34, 0.11, 0.11, 0.13, 0.31]
    ncols    = len(HEADERS)

    tbl_top  = int(MT + Inches(0.75))
    hdr_h    = int(Inches(0.50))
    row_h    = int(Inches(0.62))
    tbl_h    = hdr_h + len(rows) * row_h

    tbl = slide.shapes.add_table(len(rows) + 1, ncols, int(ML), tbl_top, int(CW), tbl_h).table
    tbl.rows[0].height = hdr_h
    for ri in range(1, len(rows) + 1):
        tbl.rows[ri].height = row_h
    for i, pct in enumerate(col_pct):
        tbl.columns[i].width = int(CW * pct)

    # Header row
    for ci, hdr in enumerate(HEADERS):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        _set_para(cell.text_frame.paragraphs[0], hdr, 10, bold=True,
                  color=WHITE, align=PP_ALIGN.CENTER)

    # Data rows: [desc, sev, prob, horizon, mitigation]
    SEV_COL, PROB_COL = 1, 2
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            if ri % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb("F5F8FF")
            else:
                cell.fill.background()
            # Color logic
            if ci == SEV_COL:
                fc = (RED if "élevé" in val.lower() or "critique" in val.lower()
                      else AMBER if "moyen" in val.lower()
                      else GREEN)
            elif ci == PROB_COL:
                fc = (RED if val == "Haute"
                      else AMBER if val == "Modérée"
                      else GREEN)
            else:
                fc = DARK
            _set_para(cell.text_frame.paragraphs[0], val, 11, color=fc,
                      align=PP_ALIGN.LEFT if ci in (0, 4) else PP_ALIGN.CENTER)
    _fit_table_rows(tbl, text_col_idx=0, font_pt=11)  # RULE 002


# ─── SLIDE 12 : PRIORITÉS (matrice impact/effort) ────────────────────────────

def _slide_priorites(prs, edm, result: dict, company: str, date_str: str, page: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header_band(slide, "PRIORITÉS", company)
    _footer_band(slide, page, company)
    _slide_title(slide, "Matrice de priorisation — Impact / Effort")

    # Reserve bottom strip for axis labels
    AXIS_H = int(Inches(0.28))
    GAP_X  = int(Inches(0.18))
    GAP_Y  = int(Inches(0.10))
    q_w    = int(CW / 2 - Inches(0.09))
    q_h    = int((CH - Inches(0.95) - AXIS_H) / 2 - GAP_Y / 2)
    q_l    = int(ML)
    q_t    = int(MT + Inches(0.68))

    quads = [
        ("IMPACT ÉLEVÉ — EFFORT FAIBLE", "Démarrer immédiatement", GREEN, _rgb("EFF8F1")),
        ("IMPACT ÉLEVÉ — EFFORT FORT",   "Planifier avec soin",    AMBER,  _rgb("FDF6EC")),
        ("IMPACT FAIBLE — EFFORT FAIBLE","Opportunités rapides",   BLUE,   _rgb("EFF5FF")),
        ("IMPACT FAIBLE — EFFORT FORT",  "Déprioritiser",          GRAY,   _rgb("F5F5F5")),
    ]

    decisions = edm.executive_decisions
    # Dynamic threshold: "high impact" = at least 15% of the biggest decision OR 10 K€ minimum.
    # mn3 fix: seuil abaissé de 50% à 15% — avec D1=119K et D2=20K, un seuil à 50%
    # (59.5K) classe D2 en "impact faible" alors que 20K€/an garanti en 1 semaine est
    # objectivement significatif. À 15% (17.85K), D2 bascule correctement en "impact élevé".
    _max_impact = max((abs(d.annual_impact or 0) for d in decisions), default=1)
    _high_impact_threshold = max(10_000, _max_impact * 0.15)
    q_items: list = [[], [], [], []]
    for dec in decisions:
        high_impact = dec.annual_impact and abs(dec.annual_impact) >= _high_impact_threshold
        dif = (dec.difficulty or "").lower()
        # Effort élevé = explicitement fort/élevé/high/difficile
        high_effort = any(w in dif for w in ("élevé", "eleve", "high", "fort", "difficile"))
        if high_impact and not high_effort:
            q_items[0].append(_sm(dec.decision))
        elif high_impact and high_effort:
            q_items[1].append(_sm(dec.decision))
        elif not high_impact and not high_effort:
            q_items[2].append(_sm(dec.decision))
        else:
            q_items[3].append(_sm(dec.decision))

    positions = [
        (q_l,                   q_t),
        (q_l + q_w + GAP_X,     q_t),
        (q_l,                   q_t + q_h + GAP_Y),
        (q_l + q_w + GAP_X,     q_t + q_h + GAP_Y),
    ]

    for idx, ((label, subtitle, c, bg), pos) in enumerate(zip(quads, positions)):
        lx, ly = pos
        _rect(slide, lx, ly, q_w, q_h, fill_color=bg, line_color=c)
        # Top accent stripe
        _rect(slide, lx, ly, q_w, int(Inches(0.04)), fill_color=c)
        _text(slide, label, lx + int(Inches(0.1)), ly + int(Inches(0.07)),
              q_w - int(Inches(0.2)), int(Inches(0.36)), size=9, bold=True, color=c)
        _text(slide, subtitle, lx + int(Inches(0.1)), ly + int(Inches(0.44)),
              q_w - int(Inches(0.2)), int(Inches(0.28)), size=9, color=GRAY, italic=True)
        items = q_items[idx] or ["—"]
        ty = ly + int(Inches(0.80))
        for it in items[:3]:
            item_h = max(int(Inches(0.32)),
                         int(_auto_row_h(it, int(q_w - Inches(0.25)), 9) / 914400 * Inches(1)))
            _text(slide, f"• {it}", lx + int(Inches(0.14)), ty,
                  q_w - int(Inches(0.24)), item_h, size=9, color=DARK)
            ty += item_h + int(Pt(3))
            if ty + item_h > ly + q_h - int(Inches(0.1)):
                break

    # ── Axis labels ──────────────────────────────────────────────────────────
    AX_Y = q_t + 2 * q_h + GAP_Y + int(Inches(0.04))
    _text(slide, "← EFFORT FAIBLE",
          q_l, AX_Y, q_w, AXIS_H,
          size=8, color=_rgb("5A6270"), italic=True)
    _text(slide, "EFFORT ÉLEVÉ →",
          q_l + q_w + GAP_X, AX_Y, q_w, AXIS_H,
          size=8, color=_rgb("5A6270"), italic=True, align=PP_ALIGN.RIGHT)

    # Spine lines — visual cross at quadrant intersection
    mid_x = int(q_l + q_w + GAP_X / 2)
    mid_y = int(q_t + q_h + GAP_Y / 2)
    _rect(slide, mid_x - 1, q_t, 2, 2 * q_h + GAP_Y, fill_color=_rgb("B0B8C4"))
    _rect(slide, q_l, mid_y - 1, 2 * q_w + GAP_X, 2, fill_color=_rgb("B0B8C4"))


# ─── SLIDE 13 : SUIVI / CARNET D'EXÉCUTION ───────────────────────────────────

def _slide_suivi(prs, edm, result: dict, company: str, date_str: str, page: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header_band(slide, "SUIVI", company)
    _footer_band(slide, page, company)
    _slide_title(slide, "Carnet d'exécution")

    exec_items = edm.execution_log or []
    decisions = edm.executive_decisions
    rows = []
    # B+C fix: le carnet utilisait exec_items (actions LLM) dont l'impact peut être négatif
    # (B: -119 000 € au lieu de +119 000 €) et le ROI calculé sans difficulté (C: 0.3 au lieu de 0.5).
    # Solution: itérer sur les phases roadmap (30→60→90j) pour le texte/dates, mais utiliser
    # les décisions EDM (par position) pour l'impact et le ROI — valeurs toujours cohérentes.
    if decisions and edm.roadmap_90_days:
        dec_list = list(decisions)
        dec_index = 0
        for phase in sorted(edm.roadmap_90_days, key=lambda p: int(str(p.horizon))):
            for action in (phase.actions or [])[:3]:
                if dec_index < len(dec_list):
                    dec = dec_list[dec_index]
                    rows.append([
                        _sm(action.decision), action.owner or dec.owner or "Direction",
                        action.due_date or "—", action.status or dec.status or "À lancer",
                        _fmt_auto(dec.annual_impact) if dec.annual_impact else "—",
                        f"{dec.roi_score:.1f}/10",
                    ])
                    dec_index += 1
                if len(rows) >= 8:
                    break
            if len(rows) >= 8:
                break
    elif exec_items:
        for item in exec_items[:8]:
            rows.append([
                _sm(item.decision), item.owner or "Direction",
                item.due_date or "—", item.status or "À lancer",
                _fmt_auto(abs(item.impact)) if item.impact else "—",
                f"{item.roi_score:.1f}/10",
            ])
    elif decisions:
        for dec in decisions[:8]:
            rows.append([
                _sm(dec.decision), dec.owner or "Direction",
                dec.timeline or "—", dec.status or "À lancer",
                _fmt_auto(dec.annual_impact) if dec.annual_impact else "—",
                f"{dec.roi_score:.1f}/10",
            ])
    if not rows:
        rows = [["—", "—", "—", "—", "—", "—"]]

    tbl_top = MT + Inches(0.55)
    hdr_h   = int(Inches(0.50))
    row_h   = int(min(Inches(0.58), (CH - Inches(0.6) - hdr_h) / max(len(rows), 1)))
    row_h   = max(int(Inches(0.40)), row_h)
    tbl_h   = hdr_h + len(rows) * row_h
    tbl = slide.shapes.add_table(len(rows) + 1, 6, ML, tbl_top, CW, tbl_h).table
    tbl.rows[0].height = hdr_h
    for _ri in range(1, len(rows) + 1):
        tbl.rows[_ri].height = row_h
    col_pct = [0.35, 0.14, 0.12, 0.13, 0.14, 0.12]
    for i, pct in enumerate(col_pct):
        tbl.columns[i].width = int(CW * pct)
    for ci, hdr in enumerate(["Action", "Responsable", "Horizon", "Statut", "Impact", "ROI"]):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        _set_para(cell.text_frame.paragraphs[0], hdr, 11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            if ri % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb("F5F8FF")
            else:
                cell.fill.background()
            fc = RED if ci == 4 else DARK
            _set_para(cell.text_frame.paragraphs[0], val, 11, color=fc,
                      align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER)
    _fit_table_rows(tbl, text_col_idx=0, font_pt=11)  # RULE 002


# ─── SLIDE 14 : PILOTAGE / TABLEAU DE BORD ───────────────────────────────────

def _slide_pilotage(prs, edm, result: dict, company: str, date_str: str, page: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header_band(slide, "PILOTAGE", company)
    _footer_band(slide, page, company)
    _slide_title(slide, "Tableau de bord de pilotage",
                 "Suivi d'avancement · Décisions prioritaires")

    decisions = edm.executive_decisions[:7]
    if not decisions:
        _text(slide, "Aucune décision à piloter.", int(ML), int(MT + Inches(0.9)),
              int(CW), int(Inches(0.5)), size=14, color=GRAY)
        return

    # ── Header bar ───────────────────────────────────────────────────────────
    HDR_Y  = int(MT + Inches(0.80))
    HDR_H  = int(Inches(0.38))
    _rect(slide, int(ML), HDR_Y, int(CW), HDR_H, fill_color=NAVY)
    col_specs = [
        ("Décision prioritaire", int(ML + Inches(0.12)), int(CW * 0.42)),
        ("Responsable",          int(ML + CW * 0.43),    int(CW * 0.13)),
        ("Impact / an",          int(ML + CW * 0.57),    int(CW * 0.13)),
        ("Horizon",              int(ML + CW * 0.71),    int(CW * 0.11)),
        ("Statut + avancement",  int(ML + CW * 0.83),    int(CW * 0.17)),
    ]
    for hdr_txt, hx, hw in col_specs:
        _text(slide, hdr_txt, hx, HDR_Y + int(Inches(0.06)), hw, int(Inches(0.26)),
              size=9, bold=True, color=WHITE)

    # ── Decision rows ─────────────────────────────────────────────────────────
    ROW_START  = HDR_Y + HDR_H + int(Inches(0.04))
    SUM_H      = int(Inches(0.36))   # bottom summary strip height
    avail_h    = int(SH - MB - ROW_START - Inches(0.04) - SUM_H - Inches(0.06))
    n          = len(decisions)
    GAP        = int(Inches(0.04))
    row_h      = max(int(Inches(0.55)), min(int(Inches(0.90)),
                     (avail_h - (n - 1) * GAP) // n))

    STATUS_COLORS = {
        "terminé":   GREEN,
        "en cours":  AMBER,
        "à lancer":  BLUE,
        "bloqué":    RED,
        "planifié":  _rgb("6B7280"),
    }

    for i, dec in enumerate(decisions):
        ry = ROW_START + i * (row_h + GAP)
        bg = _rgb("F8FAFC") if i % 2 == 0 else WHITE
        _rect(slide, int(ML), ry, int(CW), row_h, fill_color=bg)

        # Status color → left accent stripe
        status_raw = (dec.status or "À lancer").lower()
        sc = next((v for k, v in STATUS_COLORS.items() if k in status_raw), BLUE)
        _rect(slide, int(ML), ry, int(Inches(0.06)), row_h, fill_color=sc)

        # Decision title (ellipsis if long)
        title = _sm(dec.decision) or "—"
        _text(slide, title,
              int(ML + Inches(0.14)), ry + int(Inches(0.08)),
              int(CW * 0.40), row_h - int(Inches(0.16)),
              size=10, bold=True, color=DARK)

        # Owner
        owner_txt = (dec.owner or "Direction")[:22]
        _text(slide, owner_txt,
              int(ML + CW * 0.43), ry + int(Inches(0.08)),
              int(CW * 0.13), row_h - int(Inches(0.16)),
              size=9, color=_rgb("5A6270"))

        # Impact
        imp_txt = _fmt_auto(dec.annual_impact) if dec.annual_impact else "—"
        imp_col = GREEN if (dec.annual_impact or 0) >= 0 else RED
        _text(slide, imp_txt,
              int(ML + CW * 0.57), ry + int(Inches(0.08)),
              int(CW * 0.13), row_h - int(Inches(0.16)),
              size=10, bold=True, color=imp_col)

        # Horizon
        _text(slide, dec.timeline or "—",
              int(ML + CW * 0.71), ry + int(Inches(0.08)),
              int(CW * 0.11), row_h - int(Inches(0.16)),
              size=9, color=DARK)

        # Status + progress bar
        BAR_X = int(ML + CW * 0.83)
        BAR_W = int(CW * 0.16)
        _text(slide, dec.status or "À lancer",
              BAR_X, ry + int(Inches(0.06)),
              BAR_W, int(Inches(0.24)),
              size=9, bold=True, color=sc)
        # Visual bar
        pct_val = 0.0
        if "terminé" in status_raw:
            pct_val = 1.0
        elif "en cours" in status_raw:
            pct_val = 0.4
        elif "planifié" in status_raw:
            pct_val = 0.1
        bar_y = ry + int(Inches(0.31))
        bar_h = int(Inches(0.13))
        _progress_bar_shape(slide, BAR_X, bar_y, BAR_W, bar_h, pct_val, sc)

    # Bottom summary strip
    SUM_Y = ROW_START + n * (row_h + GAP) - GAP + int(Inches(0.06))
    if SUM_Y + SUM_H < int(SH - MB):
        total_imp = sum(d.annual_impact or 0 for d in decisions)
        _rect(slide, int(ML), SUM_Y, int(CW), int(Inches(0.30)), fill_color=NAVY)
        _text(slide, f"Impact total identifié : {_fmt_auto(total_imp)} / an",
              int(ML + Inches(0.2)), SUM_Y + int(Inches(0.05)),
              int(CW * 0.5), int(Inches(0.22)),
              size=10, bold=True, color=WHITE)
        _text(slide, f"{n} décisions · Révision mensuelle recommandée",
              int(ML + CW * 0.55), SUM_Y + int(Inches(0.05)),
              int(CW * 0.44), int(Inches(0.22)),
              size=9, color=_rgb("AABDD0"), align=PP_ALIGN.RIGHT)


# ─── SLIDE 15 : LUNDI MATIN ───────────────────────────────────────────────────

def _slide_lundi_matin(prs, edm, result: dict, company: str, date_str: str, page: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header_band(slide, "LUNDI MATIN", company)
    _footer_band(slide, page, company)
    _slide_title(slide, "Ce que vous faites lundi matin",
                 "3 décisions à enclencher dès cette semaine")

    top_3 = edm.executive_decisions[:3]
    card_h = (CH - Inches(0.9)) / 3 - Inches(0.1)
    card_top = MT + Inches(0.75)   # subtitle ends ~MT + 0.65"

    for i in range(3):
        dec = top_3[i] if i < len(top_3) else None
        cy = card_top + i * (card_h + Inches(0.12))
        c = [RED, AMBER, BLUE][i]
        _rect(slide, ML, cy, CW, card_h, fill_color=_rgb("F5F8FF"), line_color=c)
        _text(slide, str(i + 1), ML + Inches(0.15), cy + Inches(0.08), Inches(0.45), card_h,
              size=34, bold=True, color=c, align=PP_ALIGN.CENTER)
        dec_txt = _sm(dec.decision) if dec else "—"
        _text(slide, dec_txt, ML + Inches(0.75), cy + Inches(0.08),
              Inches(7.5), card_h - Inches(0.1), size=16, bold=True, color=DARK)
        if dec:
            items = [
                (dec.owner or "Direction", "Responsable"),
                (_fmt_auto(dec.annual_impact) if dec.annual_impact else "—", "Impact / an"),
                (dec.timeline or "—", "Horizon"),
                (f"{dec.roi_score:.1f}/10" if dec.roi_score else "—", "Score ROI"),
            ]
            for j, (val, lbl) in enumerate(items):
                lx = Inches(8.5) + j * Inches(1.0)
                _text(slide, val, lx, cy + Inches(0.08), Inches(0.95), Inches(0.55),
                      size=9, bold=True, color=c, align=PP_ALIGN.CENTER)
                _text(slide, lbl, lx, cy + Inches(0.58), Inches(0.95), Inches(0.3),
                      size=8, color=GRAY, align=PP_ALIGN.CENTER)


# ─── SLIDE 16 : ANNEXE ────────────────────────────────────────────────────────

def _slide_annexe(prs, edm, result: dict, company: str, date_str: str, page: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _header_band(slide, "ANNEXE", company)
    _footer_band(slide, page, company)
    _slide_title(slide, "Annexe — Qualité, Méthodologie et Engagements")

    dq = edm.data_quality
    score_data = dq.score_data if dq else (result.get("score_confiance") or 70)
    anomalies = (dq.anomalies if dq and dq.anomalies else result.get("coaching_issues") or [])
    assumptions = (dq.assumptions if dq and dq.assumptions else [])

    col_w = CW / 3 - Inches(0.12)
    col_top = MT + Inches(0.55)
    col_h = CH - Inches(0.6)

    # Colonne 1 : Qualité des données
    lx1 = ML
    _rect(slide, lx1, col_top, col_w, col_h, fill_color=_rgb("F5F8FF"), line_color=BLUE)
    _text(slide, "QUALITÉ DES DONNÉES", lx1 + Inches(0.1), col_top + Inches(0.1),
          col_w - Inches(0.2), Inches(0.4), size=11, bold=True, color=BLUE)
    _text(slide, f"Score qualité données : {score_data}%",
          lx1 + Inches(0.1), col_top + Inches(0.6), col_w - Inches(0.2), Inches(0.4),
          size=16, bold=True, color=NAVY)
    ay = col_top + Inches(1.1)
    if anomalies:
        for a in anomalies[:4]:
            _text(slide, f"• {str(a)}", lx1 + Inches(0.1), ay,
                  col_w - Inches(0.2), Inches(0.45), size=10, color=DARK)
            ay += Inches(0.43)
    else:
        _text(slide, "Aucune anomalie majeure détectée.", lx1 + Inches(0.1), ay,
              col_w - Inches(0.2), Inches(0.4), size=11, color=GREEN)

    # Colonne 2 : Hypothèses et Méthodologie
    lx2 = ML + col_w + Inches(0.18)
    _rect(slide, lx2, col_top, col_w, col_h, fill_color=_rgb("F5F8FF"), line_color=AMBER)
    _text(slide, "HYPOTHÈSES ET MÉTHODOLOGIE", lx2 + Inches(0.1), col_top + Inches(0.1),
          col_w - Inches(0.2), Inches(0.4), size=11, bold=True, color=AMBER)
    hyp = assumptions[:4] or [
        "Calculs basés sur les données soumises uniquement.",
        "Les projections supposent une mise en œuvre dans les 30 jours.",
        "Aucune extrapolation sectorielle n'a été effectuée.",
        "Le coût de l'inaction est calculé par interpolation linéaire.",
    ]
    hy = col_top + Inches(0.6)
    for h in hyp:
        _text(slide, f"• {str(h)}", lx2 + Inches(0.1), hy,
              col_w - Inches(0.2), Inches(0.45), size=10, color=DARK)
        hy += Inches(0.45)

    # Colonne 3 : Engagements Pepperyn
    lx3 = ML + 2 * (col_w + Inches(0.18))
    _rect(slide, lx3, col_top, col_w, col_h, fill_color=_rgb("F5F8FF"), line_color=NAVY)
    _text(slide, "ENGAGEMENTS PEPPERYN", lx3 + Inches(0.1), col_top + Inches(0.1),
          col_w - Inches(0.2), Inches(0.4), size=11, bold=True, color=NAVY)
    engagements = [
        "Accompagnement sur 90 jours pour l'implémentation.",
        "Révision bimensuelle des indicateurs de pilotage.",
        "Adaptation du plan en cas de changement de contexte.",
        "Disponibilité directe : info@finflate.com",
    ]
    ey = col_top + Inches(0.6)
    for e in engagements:
        _text(slide, f"• {e}", lx3 + Inches(0.1), ey,
              col_w - Inches(0.2), Inches(0.42), size=11, color=DARK)
        ey += Inches(0.45)

    _text(slide,
          "Ce document est confidentiel et destiné exclusivement à la direction de l'entreprise cliente. "
          "Pepperyn ne garantit pas l'exactitude des projections qui dépendent des données fournies.",
          ML, SH - Inches(0.85), CW, Inches(0.35), size=9, color=GRAY, italic=True)


# ─── POINT D'ENTRÉE ───────────────────────────────────────────────────────────

def generate_pptx_report(result: Any, company_name: Optional[str] = None) -> bytes:
    """
    Génère le Board Deck Pepperyn (16 slides).

    Args:
        result: ExecutiveCaseJSON (V2 — source unique de vérité)
                OU dict brut de l'analyse LLM (legacy — rétrocompatibilité).
        company_name: nom de la société (couverture).

    Pipeline V2 :
        - Si ExecutiveCaseJSON → adapté via case_to_edm() + case_to_result_dict()
          Garantie : PDF, PPTX et Excel affichent exactement les mêmes chiffres.
        - Si dict → comportement legacy inchangé.
    """
    if isinstance(result, ExecutiveCaseJSON):
        # ── V2 : source unique de vérité ─────────────────────────────────────
        from services.executive_case_builder import case_to_edm, case_to_result_dict
        edm        = case_to_edm(result)
        result_raw = case_to_result_dict(result)
        company    = company_name or result.company_name or "—"
        date_str   = result.analysis_date or _fr_date(datetime.now())
        result     = result_raw   # alias pour les builders internes (inchangés)
    else:
        # ── Legacy : comportement existant ───────────────────────────────────
        edm      = build_executive_decision_model(result)
        date_str = _fr_date(datetime.now())
        company  = company_name or result.get("company_name") or "—"

    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    slides_builders = [
        _slide_cover,                        # S1
        _slide_exec_summary,                 # S2
        _slide_diagnostic,                   # S3
        _slide_dashboard,                    # S4
        _slide_impact_financier,             # S5
        _slide_cout_inaction,                # S6
        _slide_decisions_prioritaires,       # S7
        _slide_raisonnement_comparatif,      # S7b — EDX-002 (omit si pas de données)
        _slide_execution,                    # S8
        _slide_simulation,                   # S9
        _slide_bridge_historique,            # S10 — Bridge passé : Normatif → Actuel
        _slide_bridge_financier,             # S11 — Bridge futur  : Actuel → Cible
        _slide_projection,                   # S12
        _slide_creation_valeur,              # S12 — Création de valeur (Phase 1 + Phase 2)
        _slide_risques,                      # S12
        _slide_priorites,                    # S12
        _slide_suivi,                        # S13
        _slide_pilotage,                     # S14
        _slide_lundi_matin,                  # S15
        _slide_annexe,                       # S16
    ]

    # RULE — page_num tracks actual slides added (not builder index).
    # If a builder adds no slide (early return), the counter does not increment,
    # so subsequent slides keep correct footers.
    actual_page = 0
    for builder in slides_builders:
        n_before = len(prs.slides)
        try:
            builder(prs, edm, result, company, date_str, actual_page + 1)
        except Exception as exc:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            _header_band(slide, f"SLIDE {actual_page + 1}", company)
            _footer_band(slide, actual_page + 1, company)
            _text(slide, f"Slide {actual_page + 1} — Données indisponibles",
                  ML, MT + Inches(1), CW, Inches(1), size=18, color=GRAY)
        if len(prs.slides) > n_before:
            actual_page += 1

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
